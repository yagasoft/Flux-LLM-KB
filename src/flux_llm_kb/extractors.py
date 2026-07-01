from __future__ import annotations

import base64
import bz2
import binascii
import csv
from dataclasses import dataclass, field, replace
from email import policy as email_policy
from email.parser import BytesParser
import gzip
import hashlib
from html import unescape
import importlib.util
from io import BytesIO
import json
import lzma
import mailbox
import math
import mimetypes
import os
from pathlib import Path
from pathlib import PurePosixPath
import re
import shutil
import sqlite3
import struct
import subprocess
import tarfile
import tempfile
from typing import Any, Callable, Iterable
from urllib.error import HTTPError, URLError
from urllib.parse import quote, unquote
from urllib.request import Request, urlopen
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile
import zlib

from .acceleration import resolve_cache_layout, validate_local_model_base_url
from .code_index import parse_code_file, references_to_metadata, symbols_to_metadata
from .crawler import (
    ARCHIVE_COMPOUND_SUFFIXES,
    ARCHIVE_EXTENSIONS,
    AUDIO_EXTENSIONS,
    AssetChunk,
    CAD_EXTENSIONS,
    CALENDAR_EXTENSIONS,
    CODE_EXTENSIONS,
    CONTAINER_EXTENSIONS,
    CONTACT_EXTENSIONS,
    DIAGRAM_COMPOUND_SUFFIXES,
    DIAGRAM_EXTENSIONS,
    DOCUMENT_EXTENSIONS,
    GEOSPATIAL_EXTENSIONS,
    IMAGE_EXTENSIONS,
    MAIL_EXTENSIONS,
    REPORT_EXTENSIONS,
    REPORT_NAMES,
    SCIENTIFIC_EXTENSIONS,
    SENSITIVE_METADATA_EXTENSIONS,
    STRUCTURED_DATA_EXTENSIONS,
    SUBTITLE_EXTENSIONS,
    TEXT_EXTENSIONS,
    VIDEO_EXTENSIONS,
    CorpusPolicy,
    classify_file,
)
from .processes import run_no_window
from .redaction import redact_text
from .text_safety import decode_text_bytes, read_text_with_bom


DIAGRAM_MAX_ZIP_MEMBERS = 1024
DIAGRAM_MAX_PAGE_XML_MEMBERS = 200
DIAGRAM_MAX_TOTAL_BYTES = 25 * 1024 * 1024
DIAGRAM_MAX_MEMBER_BYTES = 5 * 1024 * 1024

LEGACY_WORD_EXTENSIONS = {".doc", ".dot", ".rtf"}
CONVERTED_WORD_EXTENSIONS = {".docm", ".dotm", ".dotx", ".odt", ".ott"}
OPENPYXL_EXTENSIONS = {".xlsx", ".xlsm", ".xltm", ".xltx"}
LEGACY_SPREADSHEET_EXTENSIONS = {".xls", ".xlsb", ".xlt"}
OPENDOCUMENT_SPREADSHEET_EXTENSIONS = {".ods", ".ots"}
PPTX_PACKAGE_EXTENSIONS = {".pptx", ".pptm", ".potm", ".potx", ".ppsm", ".ppsx"}
LEGACY_PRESENTATION_EXTENSIONS = {".ppt", ".pot", ".pps"}
OPENDOCUMENT_PRESENTATION_EXTENSIONS = {".odp", ".otp"}
LOCAL_PUBLICATION_EXTENSIONS = {".epub", ".fb2"}
CALIBRE_PUBLICATION_EXTENSIONS = {".mobi", ".azw", ".azw3", ".lit"}
COMIC_ARCHIVE_EXTENSIONS = {".cbz", ".cbr", ".cb7", ".cbt"}
MEDIA_TRANSCRIPT_SIDECAR_SUFFIXES = (".txt", ".md", ".vtt", ".srt")
OCR_CACHE_SCHEMA = "flux-ocr-cache-v1"
ASR_CACHE_SCHEMA = "flux-asr-cache-v1"
VISION_CACHE_SCHEMA = "flux-vision-cache-v1"
THUMBNAIL_CACHE_SCHEMA = "flux-thumbnail-cache-v1"
VISION_PROMPT_SCHEMA = "flux-vision-caption-v3"
OCR_MAX_PDF_PAGES = 25
OCR_PDF_PAGE_BATCH_SIZE = 5
OCR_PDF_DPI = 200
OCR_MAX_IMAGE_EDGE = 6000
OCR_TIMEOUT_SECONDS = 30
ASR_AUDIO_SAMPLE_RATE = 16000
ASR_FFMPEG_TIMEOUT_SECONDS = 300
MEDIA_ASR_SEGMENT_SECONDS = 15 * 60
MEDIA_SEGMENT_CHUNK_INDEX_STRIDE = 1000
VIDEO_FRAME_CHUNK_INDEX_BASE = 100_000
PDF_OCR_CHUNK_INDEX_BASE = 200_000
VISION_TIMEOUT_SECONDS = 180
VISION_REQUEST_MAX_EDGE = 1280
VISION_NUM_PREDICT = 1024
FRAME_SAMPLING_TIMEOUT_SECONDS = 120
PRACTICAL_TEXT_LIMIT_BYTES = 2 * 1024 * 1024
PRACTICAL_SUMMARY_LIMIT = 50
DECORATIVE_ICON_MAX_DIMENSION = 32
DECORATIVE_ICON_MAX_BYTES = 4096


@dataclass(frozen=True)
class ContainerChildAsset:
    member_path: str
    file_kind: str
    mime_type: str | None
    extension: str
    size_bytes: int
    quick_hash: str
    content_hash: str | None
    extraction_tier: str
    extraction_status: str
    chunks: tuple[AssetChunk, ...] = ()
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class ExtractionResult:
    status: str
    chunks: tuple[AssetChunk, ...] = ()
    child_assets: tuple[ContainerChildAsset, ...] = ()
    metadata: dict[str, Any] = field(default_factory=dict)
    message: str | None = None


@dataclass(frozen=True)
class OcrResult:
    status: str
    text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    message: str | None = None


@dataclass(frozen=True)
class AsrResult:
    status: str
    text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    message: str | None = None


@dataclass(frozen=True)
class VisionResult:
    status: str
    text: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)
    message: str | None = None


@dataclass(frozen=True)
class FrameSamplingResult:
    status: str
    chunks: tuple[AssetChunk, ...] = ()
    metadata: dict[str, Any] = field(default_factory=dict)
    message: str | None = None


def _is_vss_shadow_path(path: str | Path) -> bool:
    text = str(path).replace("/", "\\").lower()
    return text.startswith("\\\\?\\globalroot\\device\\harddiskvolumeshadowcopy")


def _extractor_path(path: str | Path) -> Path:
    file_path = Path(path).expanduser()
    if _is_vss_shadow_path(file_path):
        return file_path
    return file_path.resolve()


def extract_file(path: str | Path, policy: CorpusPolicy, *, relative_path: str | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    classification = classify_file(file_path, policy)
    sample_first = _extract_sample_first_structured(file_path, policy)
    if sample_first is not None:
        return sample_first
    if classification.file_kind in {"text", "code"}:
        return _extract_text(file_path, policy, extractor=classification.file_kind, relative_path=relative_path)
    if classification.file_kind == "subtitle":
        return _extract_subtitle(file_path)
    if classification.file_kind == "mail":
        return _extract_mail(file_path)
    if classification.file_kind == "calendar":
        return _extract_calendar(file_path)
    if classification.file_kind == "contact":
        return _extract_contact(file_path)
    if classification.file_kind == "structured_data":
        return _extract_structured_data(file_path, policy)
    if classification.file_kind == "report":
        return _extract_report(file_path)
    if classification.file_kind == "database":
        return _extract_database(file_path)
    if classification.file_kind in {"geospatial", "cad", "scientific"}:
        return _extract_domain_metadata(file_path, classification.file_kind)
    if classification.file_kind == "sensitive_metadata":
        return _extract_sensitive_metadata(file_path)
    if classification.file_kind == "document":
        return _extract_document(file_path, policy)
    if classification.file_kind in {"archive", "container"}:
        return _extract_container(file_path, policy, container_kind=classification.file_kind)
    if classification.file_kind == "diagram":
        return _extract_diagram(file_path)
    if classification.file_kind == "image":
        return _extract_image(file_path)
    if classification.file_kind in {"audio", "video"}:
        return _extract_media(file_path, classification.file_kind)
    return ExtractionResult(
        status="metadata_only",
        metadata={"extractor": classification.file_kind, "mime_type": classification.mime_type},
    )


def extractor_availability() -> dict[str, dict[str, Any]]:
    return {
        "pypdf": _module_check("pypdf"),
        "python_docx": _module_check("docx"),
        "python_pptx": _module_check("pptx"),
        "openpyxl": _module_check("openpyxl"),
        "pillow": _module_check("PIL"),
        "watchdog": _module_check("watchdog"),
        "libreoffice": _first_tool_check("LibreOffice", ("soffice", "libreoffice")),
        "antiword": _tool_check("antiword"),
        "catdoc": _tool_check("catdoc"),
        "wvText": _tool_check("wvText"),
        "word_com": _word_com_check(),
        "excel_com": _excel_com_check(),
        "powerpoint_com": _powerpoint_com_check(),
        "pdftoppm": _tool_check("pdftoppm"),
        "seven_zip": _first_tool_check("7-Zip", ("7z", "7zz")),
        "bsdtar": _tool_check("bsdtar"),
        "unar": _tool_check("unar"),
        "unrar": _tool_check("unrar"),
        "zstd": _tool_check("zstd"),
        "lz4": _tool_check("lz4"),
        "ar": _tool_check("ar"),
        "rpm2cpio": _tool_check("rpm2cpio"),
        "ffprobe": _tool_check("ffprobe"),
        "ffmpeg": _tool_check("ffmpeg"),
        "tesseract": _tool_check("tesseract"),
        "faster_whisper": _module_check("faster_whisper"),
        "ebook_convert": _tool_check("ebook-convert"),
        "readpst": _tool_check("readpst"),
        "msgconvert": _tool_check("msgconvert"),
        "duckdb": _module_check("duckdb"),
        "pyarrow": _module_check("pyarrow"),
        "ogrinfo": _tool_check("ogrinfo"),
        "gdalinfo": _tool_check("gdalinfo"),
        "ifcopenshell": _module_check("ifcopenshell"),
        "assimp": _tool_check("assimp"),
        "blender": _tool_check("blender"),
        "exiftool": _tool_check("exiftool"),
        "pandoc": _tool_check("pandoc"),
    }


def image_metadata(path: str | Path) -> dict[str, Any]:
    file_path = Path(path)
    metadata: dict[str, Any] = {"extractor": "image"}
    try:
        metadata["size_bytes"] = file_path.stat().st_size
    except OSError:
        pass
    dimensions = _image_dimensions(file_path)
    if dimensions:
        metadata.update({"width": dimensions[0], "height": dimensions[1]})
    has_transparency = _image_has_transparency(file_path)
    if has_transparency is not None:
        metadata["has_transparency"] = has_transparency
    return metadata


def _extract_text(path: Path, policy: CorpusPolicy, *, extractor: str, relative_path: str | None = None) -> ExtractionResult:
    if path.stat().st_size > policy.max_inline_bytes:
        return ExtractionResult(
            status="blocked_missing_dependency",
            metadata={"extractor": extractor},
            message="text file exceeds inline extraction limit",
        )
    if extractor == "code":
        parsed = parse_code_file(path, root=policy.root_path, relative_path=relative_path)
        symbol_metadata = symbols_to_metadata(parsed.symbols)
        reference_metadata = references_to_metadata(parsed.references)
        chunks = tuple(
            AssetChunk(
                chunk_index=chunk.chunk_index,
                title=chunk.title,
                body=chunk.body,
                modality="code",
                locator=chunk.locator,
                token_estimate=chunk.token_estimate,
                metadata={
                    **chunk.metadata,
                    "code_symbols": [symbol for symbol in symbol_metadata if symbol.get("chunk_index") == chunk.chunk_index],
                    "code_references": [reference for reference in reference_metadata if reference.get("chunk_index") == chunk.chunk_index],
                },
            )
            for chunk in parsed.chunks
        )
        metadata = {
            "extractor": extractor,
            "code": {
                **parsed.metadata,
                "symbols": symbol_metadata,
                "references": reference_metadata,
            },
        }
        return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata=metadata)
    text = read_text_with_bom(path).strip()
    chunks = _chunks_from_text(text, path.name)
    metadata: dict[str, Any] = {"extractor": extractor}
    if not chunks and path.stat().st_size == 0:
        metadata["empty"] = True
        return ExtractionResult(status="indexed", chunks=chunks, metadata=metadata)
    return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata=metadata)


def _extract_sample_first_structured(path: Path, policy: CorpusPolicy) -> ExtractionResult | None:
    if path.stat().st_size <= policy.max_inline_bytes:
        return None
    ext = path.suffix.lower()
    if ext in {".csv", ".tsv", ".psv", ".ssv"}:
        delimiters = {".csv": ",", ".tsv": "\t", ".psv": "|", ".ssv": " "}
        return _extract_sample_first_delimited(path, delimiter=delimiters[ext], source_format=ext.lstrip("."))
    if ext == ".json":
        return _extract_sample_first_json(path)
    if ext in {".jsonl", ".ndjson"}:
        return _extract_sample_first_jsonl(path, source_format=ext.lstrip("."))
    if ext == ".jsonld":
        return _extract_sample_first_json(path, source_format="jsonld")
    if ext in OPENPYXL_EXTENSIONS:
        return _extract_sample_first_workbook(path, extractor=ext.lstrip("."))
    return None


def _extract_sample_first_delimited(
    path: Path,
    *,
    delimiter: str,
    source_format: str | None = None,
    sample_limit: int = 10,
) -> ExtractionResult:
    rows: list[dict[str, str]] = []
    row_count = 0
    columns: list[str] = []
    format_name = source_format or ("tsv" if delimiter == "\t" else "csv")
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.DictReader(handle, delimiter=delimiter)
        columns = list(reader.fieldnames or [])
        for row in reader:
            row_count += 1
            if len(rows) < sample_limit:
                rows.append({str(key): str(value) for key, value in row.items() if key is not None})
    sample = {
        "format": format_name,
        "columns": columns,
        "row_count_estimate": row_count,
        "sample_row_count": len(rows),
        "truncated": row_count > len(rows),
        "parse_status": "sampled",
    }
    body = _sample_first_body(title=path.name, sample=sample, rows=rows)
    chunk = AssetChunk(
        chunk_index=0,
        title=f"{path.name} sample",
        body=body,
        token_estimate=max(1, len(body) // 4),
        metadata={"sample_first": True, "format": sample["format"], "columns": columns},
    )
    return ExtractionResult(status="indexed", chunks=(chunk,), metadata={"extractor": "sample_first_tabular", "sample_first": sample})


def _extract_sample_first_jsonl(path: Path, *, source_format: str = "jsonl", sample_limit: int = 10) -> ExtractionResult:
    rows: list[dict[str, Any]] = []
    columns: set[str] = set()
    row_count = 0
    parse_errors = 0
    with path.open("r", encoding="utf-8", errors="replace") as handle:
        for line in handle:
            stripped = line.strip()
            if not stripped:
                continue
            row_count += 1
            try:
                payload = json.loads(stripped)
            except json.JSONDecodeError:
                parse_errors += 1
                continue
            if isinstance(payload, dict):
                columns.update(str(key) for key in payload)
                if len(rows) < sample_limit:
                    rows.append(payload)
    ordered_columns = sorted(columns)
    sample = {
        "format": source_format,
        "columns": ordered_columns,
        "row_count_estimate": row_count,
        "sample_row_count": len(rows),
        "truncated": row_count > len(rows),
        "parse_status": "sampled_with_errors" if parse_errors else "sampled",
        "parse_errors": parse_errors,
    }
    body = _sample_first_body(title=path.name, sample=sample, rows=rows)
    chunk = AssetChunk(
        chunk_index=0,
        title=f"{path.name} sample",
        body=body,
        token_estimate=max(1, len(body) // 4),
        metadata={"sample_first": True, "format": source_format, "columns": ordered_columns},
    )
    return ExtractionResult(status="indexed", chunks=(chunk,), metadata={"extractor": "sample_first_jsonl", "sample_first": sample})


def _extract_sample_first_json(path: Path, *, source_format: str = "json", sample_limit: int = 10) -> ExtractionResult:
    try:
        payload = json.loads(read_text_with_bom(path))
    except json.JSONDecodeError as exc:
        return ExtractionResult(
            status="blocked_missing_dependency",
            metadata={"extractor": "sample_first_json", "sample_first": {"format": source_format, "parse_status": "parse_error"}},
            message=f"JSON parse failed: {exc.msg}",
        )
    rows: list[dict[str, Any]]
    source_key: str | None = None
    if isinstance(payload, list):
        source_rows = payload
    elif isinstance(payload, dict):
        source_rows = None
        for key, value in payload.items():
            if isinstance(value, list):
                source_key = str(key)
                source_rows = value
                break
        if source_rows is None:
            source_rows = [payload]
    else:
        source_rows = [{"value": payload}]
    rows = [_coerce_sample_row(row) for row in source_rows[:sample_limit]]
    columns = sorted({key for row in source_rows if isinstance(row, dict) for key in row})
    if not columns and rows:
        columns = sorted({key for row in rows for key in row})
    sample = {
        "format": source_format,
        "columns": [str(column) for column in columns],
        "row_count_estimate": len(source_rows),
        "sample_row_count": len(rows),
        "truncated": len(source_rows) > len(rows),
        "parse_status": "sampled",
    }
    if source_key is not None:
        sample["source_key"] = source_key
    body = _sample_first_body(title=path.name, sample=sample, rows=rows)
    chunk = AssetChunk(
        chunk_index=0,
        title=f"{path.name} sample",
        body=body,
        token_estimate=max(1, len(body) // 4),
        metadata={"sample_first": True, "format": source_format, "columns": sample["columns"]},
    )
    return ExtractionResult(status="indexed", chunks=(chunk,), metadata={"extractor": "sample_first_json", "sample_first": sample})


def _extract_sample_first_workbook(path: Path, *, extractor: str, sample_limit: int = 10) -> ExtractionResult:
    try:
        import openpyxl
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": extractor}, message="openpyxl not installed")
    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        if _is_invalid_package_error(exc):
            return _invalid_package_result(exc, extractor=extractor)
        raise
    rows: list[dict[str, Any]] = []
    row_count = 0
    columns: list[str] = []
    sheet_names: list[str] = []
    try:
        for worksheet in workbook.worksheets[:10]:
            sheet_names.append(str(worksheet.title))
            sheet_columns: list[str] | None = None
            for raw_row in worksheet.iter_rows(values_only=True):
                values = list(raw_row[:30])
                if not any(value is not None for value in values):
                    continue
                if sheet_columns is None:
                    sheet_columns = [str(value) if value is not None else f"column_{index + 1}" for index, value in enumerate(values)]
                    if not columns:
                        columns = sheet_columns
                    continue
                row_count += 1
                if len(rows) < sample_limit:
                    rows.append(
                        {
                            str(column): _sample_value(value)
                            for column, value in zip(sheet_columns, values, strict=False)
                            if value is not None
                        }
                        | {"_sheet": str(worksheet.title)}
                    )
    finally:
        close = getattr(workbook, "close", None)
        if callable(close):
            close()
    sample = {
        "format": "workbook",
        "columns": columns,
        "row_count_estimate": row_count,
        "sample_row_count": len(rows),
        "truncated": row_count > len(rows),
        "parse_status": "sampled",
        "sheet_count": len(getattr(workbook, "worksheets", [])),
        "sheet_names": sheet_names,
        "source_extension": path.suffix.lower(),
    }
    body = _sample_first_body(title=path.name, sample=sample, rows=rows)
    chunk = AssetChunk(
        chunk_index=0,
        title=f"{path.name} sample",
        body=body,
        token_estimate=max(1, len(body) // 4),
        metadata={"sample_first": True, "format": "workbook", "columns": columns},
    )
    return ExtractionResult(status="indexed", chunks=(chunk,), metadata={"extractor": "sample_first_workbook", "sample_first": sample})


def _coerce_sample_row(row: Any) -> dict[str, Any]:
    if isinstance(row, dict):
        return {str(key): _sample_value(value) for key, value in row.items()}
    return {"value": _sample_value(row)}


def _sample_value(value: Any) -> Any:
    if value is None or isinstance(value, (bool, int, float, str)):
        return value
    return str(value)


def _sample_first_body(*, title: str, sample: dict[str, Any], rows: list[dict[str, Any]]) -> str:
    return "\n".join(
        [
            f"# {title} sample-first profile",
            f"Format: {sample.get('format')}",
            f"Rows estimated: {sample.get('row_count_estimate')}",
            f"Columns: {', '.join(sample.get('columns') or [])}",
            "Sample rows:",
            json.dumps(rows, ensure_ascii=True, sort_keys=True, default=str),
        ]
    )


def _extract_subtitle(path: Path) -> ExtractionResult:
    text = _read_text_limited(path)
    lines: list[str] = []
    cue_count = 0
    in_note_block = False
    for raw_line in text.splitlines():
        line = raw_line.strip("\ufeff").strip()
        upper = line.upper()
        if not line:
            in_note_block = False
            continue
        if upper.startswith("NOTE"):
            in_note_block = True
            continue
        if in_note_block:
            continue
        if upper in {"WEBVTT", "STYLE", "REGION"} or line.startswith("["):
            continue
        if upper.startswith("FORMAT:"):
            continue
        if re.fullmatch(r"\d+", line):
            continue
        if "-->" in line:
            cue_count += 1
            continue
        if upper.startswith("DIALOGUE:"):
            cue_count += 1
            line = _ass_dialogue_text(line)
        cleaned = _clean_subtitle_line(line)
        if cleaned:
            lines.append(cleaned)
    body = "\n".join(lines)
    chunks = _chunks_from_text(body, path.name, modality="transcript")
    if cue_count == 0:
        cue_count = len(lines)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "subtitle",
            "subtitle_format": path.suffix.lower().lstrip(".") or "unknown",
            "cue_count": cue_count,
        },
    )


def _extract_mail(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    if ext == ".msg":
        return _extract_msg(path)
    messages = _mail_messages(path) if ext == ".mbox" else [_parse_email_message(path.read_bytes())]
    return _mail_result(messages, path.name, mail_format=ext.lstrip(".") or "unknown")


def _extract_msg(path: Path) -> ExtractionResult:
    msgconvert = shutil.which("msgconvert")
    if not msgconvert:
        return ExtractionResult(
            status="blocked_missing_dependency",
            metadata={"extractor": "mail", "mail_format": "msg", "dependency": "msgconvert"},
            message="msgconvert command not found",
        )
    with tempfile.TemporaryDirectory(prefix="flux-kb-msg-") as temp_dir:
        output_path = Path(temp_dir) / f"{path.stem}.eml"
        result = run_no_window([msgconvert, "--outfile", str(output_path), str(path)], timeout=OCR_TIMEOUT_SECONDS)
        if result.returncode != 0 or not output_path.exists():
            message = (result.stderr or result.stdout or "msgconvert failed").strip()
            return ExtractionResult(
                status="blocked_missing_dependency",
                metadata={"extractor": "mail", "mail_format": "msg", "dependency": "msgconvert"},
                message=message,
            )
        converted = _parse_email_message(output_path.read_bytes())
    mail_result = _mail_result([converted], path.name, mail_format="msg")
    return ExtractionResult(
        status=mail_result.status,
        chunks=mail_result.chunks,
        child_assets=mail_result.child_assets,
        metadata={**mail_result.metadata, "converted_format": "eml"},
        message=mail_result.message,
    )


def _mail_result(messages: list[Any], source_name: str, *, mail_format: str) -> ExtractionResult:
    parts: list[str] = []
    subjects: list[str] = []
    attachment_count = 0
    for index, message in enumerate(messages[:PRACTICAL_SUMMARY_LIMIT], start=1):
        subject = str(message.get("Subject") or "").strip()
        if subject:
            subjects.append(subject)
        attachment_count += _mail_attachment_count(message)
        body = _mail_plain_text(message)
        message_parts = [f"Message {index}"]
        if subject:
            message_parts.append(f"Subject: {subject}")
        if body:
            message_parts.append(body)
        parts.append("\n".join(message_parts))
    text = "\n\n".join(part for part in parts if part.strip())
    chunks = _chunks_from_text(text, source_name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "mail",
            "mail_format": mail_format,
            "message_count": len(messages),
            "attachment_count": attachment_count,
            "subjects": subjects[:20],
            "truncated": len(messages) > PRACTICAL_SUMMARY_LIMIT,
        },
    )


def _extract_calendar(path: Path) -> ExtractionResult:
    blocks = _ical_blocks(_read_text_limited(path), "VEVENT")
    parts: list[str] = []
    for index, block in enumerate(blocks[:PRACTICAL_SUMMARY_LIMIT], start=1):
        lines = [f"Event {index}"]
        for label, key in (("Summary", "SUMMARY"), ("Start", "DTSTART"), ("End", "DTEND"), ("Location", "LOCATION"), ("Description", "DESCRIPTION")):
            value = _first_property_value(block, key)
            if value:
                lines.append(f"{label}: {value}")
        parts.append("\n".join(lines))
    chunks = _chunks_from_text("\n\n".join(parts), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "calendar",
            "calendar_format": path.suffix.lower().lstrip(".") or "unknown",
            "event_count": len(blocks),
            "truncated": len(blocks) > PRACTICAL_SUMMARY_LIMIT,
        },
    )


def _extract_contact(path: Path) -> ExtractionResult:
    blocks = _ical_blocks(_read_text_limited(path), "VCARD")
    parts: list[str] = []
    for index, block in enumerate(blocks[:PRACTICAL_SUMMARY_LIMIT], start=1):
        lines = [f"Contact {index}"]
        for label, key in (("Name", "FN"), ("Organization", "ORG"), ("Title", "TITLE"), ("Note", "NOTE")):
            value = _first_property_value(block, key)
            if value:
                lines.append(f"{label}: {value}")
        parts.append("\n".join(lines))
    chunks = _chunks_from_text("\n\n".join(parts), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "contact",
            "contact_format": path.suffix.lower().lstrip(".") or "unknown",
            "contact_count": len(blocks),
            "truncated": len(blocks) > PRACTICAL_SUMMARY_LIMIT,
        },
    )


def _extract_structured_data(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    ext = path.suffix.lower()
    if ext in {".csv", ".tsv", ".psv", ".ssv", ".json", ".jsonl", ".ndjson", ".jsonld"}:
        expanded_policy = replace(policy, max_inline_bytes=-1)
        sample = _extract_sample_first_structured(path, expanded_policy)
        if sample is not None:
            return sample
    if path.stat().st_size <= PRACTICAL_TEXT_LIMIT_BYTES:
        chunks = _chunks_from_text(read_text_with_bom(path), path.name)
        return ExtractionResult(
            status="indexed" if chunks else "metadata_only",
            chunks=chunks,
            metadata={"extractor": "structured_data", "source_extension": ext},
        )
    return ExtractionResult(
        status="metadata_only",
        metadata={"extractor": "structured_data", "source_extension": ext, "reason": "size_limit"},
    )


def _extract_report(path: Path) -> ExtractionResult:
    name = path.name.lower()
    ext = path.suffix.lower()
    if ext == ".sarif":
        return _extract_sarif_report(path)
    if ext == ".cyclonedx":
        return _extract_cyclonedx_report(path)
    if ext == ".spdx":
        return _extract_spdx_report(path)
    if ext == ".har":
        return _extract_har_report(path)
    if ext == ".lcov":
        return _extract_lcov_report(path)
    if ext == ".tap":
        return _extract_tap_report(path)
    if ext == ".trx":
        return _extract_trx_report(path)
    if ext == ".xml" or name in REPORT_NAMES:
        return _extract_xml_report(path)
    return _extract_bounded_report_text(path)


def _extract_sarif_report(path: Path) -> ExtractionResult:
    payload = json.loads(_read_text_limited(path))
    runs = payload.get("runs") if isinstance(payload, dict) else None
    lines: list[str] = []
    finding_count = 0
    if isinstance(runs, list):
        for run in runs:
            if not isinstance(run, dict):
                continue
            tool = (((run.get("tool") or {}).get("driver") or {}).get("name") if isinstance(run.get("tool"), dict) else None)
            results = run.get("results")
            if not isinstance(results, list):
                continue
            finding_count += sum(1 for result in results if isinstance(result, dict))
            if tool:
                lines.append(f"Tool: {tool}")
            for result in results[:PRACTICAL_SUMMARY_LIMIT]:
                if not isinstance(result, dict):
                    continue
                rule_id = str(result.get("ruleId") or "finding")
                message = result.get("message")
                message_text = str((message or {}).get("text") or (message or {}).get("markdown") or "") if isinstance(message, dict) else ""
                lines.append(f"{rule_id}: {message_text}".strip())
    chunks = _chunks_from_text("\n".join(lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": "report", "report_format": "sarif", "finding_count": finding_count},
    )


def _extract_cyclonedx_report(path: Path) -> ExtractionResult:
    payload = json.loads(_read_text_limited(path))
    components = payload.get("components") if isinstance(payload, dict) else None
    lines: list[str] = []
    if isinstance(components, list):
        for component in components[:PRACTICAL_SUMMARY_LIMIT]:
            if not isinstance(component, dict):
                continue
            name = str(component.get("name") or "").strip()
            version = str(component.get("version") or "").strip()
            purl = str(component.get("purl") or "").strip()
            label = " ".join(part for part in (name, version) if part)
            if purl:
                label = f"{label} {purl}".strip()
            if label:
                lines.append(label)
    chunks = _chunks_from_text("\n".join(lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "cyclonedx",
            "component_count": len(components) if isinstance(components, list) else 0,
        },
    )


def _extract_spdx_report(path: Path) -> ExtractionResult:
    packages: list[dict[str, str]] = []
    current: dict[str, str] = {}
    for raw_line in _read_text_limited(path).splitlines():
        line = raw_line.strip()
        if not line or ":" not in line:
            continue
        key, _, value = line.partition(":")
        key = key.strip()
        value = value.strip()
        if key == "PackageName":
            if current:
                packages.append(current)
            current = {"name": value}
        elif key == "PackageVersion" and current:
            current["version"] = value
    if current:
        packages.append(current)
    lines = [
        " ".join(part for part in (package.get("name"), package.get("version")) if part)
        for package in packages[:PRACTICAL_SUMMARY_LIMIT]
    ]
    chunks = _chunks_from_text("\n".join(lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": "report", "report_format": "spdx", "package_count": len(packages)},
    )


def _extract_har_report(path: Path) -> ExtractionResult:
    payload = json.loads(_read_text_limited(path))
    entries = (((payload or {}).get("log") or {}).get("entries") if isinstance(payload, dict) else None) or []
    lines: list[str] = []
    if isinstance(entries, list):
        for entry in entries[:PRACTICAL_SUMMARY_LIMIT]:
            if not isinstance(entry, dict):
                continue
            request = entry.get("request") if isinstance(entry.get("request"), dict) else {}
            response = entry.get("response") if isinstance(entry.get("response"), dict) else {}
            method = str(request.get("method") or "GET")
            url = str(request.get("url") or "")
            status = response.get("status")
            if url:
                lines.append(f"{method} {url} -> {status}")
    chunks = _chunks_from_text("\n".join(lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": "report", "report_format": "har", "entry_count": len(entries) if isinstance(entries, list) else 0},
    )


def _extract_xml_report(path: Path) -> ExtractionResult:
    root = ElementTree.fromstring(_read_text_limited(path))
    local_name = _xml_local_name(root.tag).lower()
    if local_name == "coverage":
        return _extract_coverage_xml_report(path, root=root)
    if local_name == "testrun":
        return _extract_trx_report(path, root=root)
    return _extract_junit_report(path, root=root)


def _extract_junit_report(path: Path, *, root: ElementTree.Element | None = None) -> ExtractionResult:
    root = root if root is not None else ElementTree.fromstring(_read_text_limited(path))
    suites = [root] if _xml_local_name(root.tag) == "testsuite" else [item for item in root.iter() if _xml_local_name(item.tag) == "testsuite"]
    test_count = sum(_int_or_none(suite.get("tests")) or 0 for suite in suites)
    failure_count = sum(_int_or_none(suite.get("failures")) or 0 for suite in suites)
    error_count = sum(_int_or_none(suite.get("errors")) or 0 for suite in suites)
    skipped_count = sum(_int_or_none(suite.get("skipped")) or 0 for suite in suites)
    if test_count == 0:
        test_count = len([item for item in root.iter() if _xml_local_name(item.tag) == "testcase"])
    failed_cases: list[str] = []
    for testcase in root.iter():
        if _xml_local_name(testcase.tag) != "testcase":
            continue
        has_failure = any(_xml_local_name(child.tag) in {"failure", "error"} for child in testcase)
        if has_failure:
            case_name = ".".join(part for part in (testcase.get("classname"), testcase.get("name")) if part)
            failed_cases.append(case_name or "failed testcase")
    body_lines = [
        f"Tests: {test_count}",
        f"Failures: {failure_count}",
        f"Errors: {error_count}",
        f"Skipped: {skipped_count}",
    ]
    if failed_cases:
        body_lines.append("Failed cases: " + ", ".join(failed_cases[:PRACTICAL_SUMMARY_LIMIT]))
    chunks = _chunks_from_text("\n".join(body_lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "junit",
            "test_count": test_count,
            "failure_count": failure_count,
            "error_count": error_count,
            "skipped_count": skipped_count,
        },
    )


def _extract_trx_report(path: Path, *, root: ElementTree.Element | None = None) -> ExtractionResult:
    root = root if root is not None else ElementTree.fromstring(_read_text_limited(path))
    results = [item for item in root.iter() if _xml_local_name(item.tag).lower() == "unittestresult"]
    failed = [
        str(item.get("testName") or item.get("testId") or "failed test")
        for item in results
        if str(item.get("outcome") or "").lower() in {"failed", "error", "timeout", "aborted"}
    ]
    body_lines = [f"Tests: {len(results)}", f"Failed: {len(failed)}"]
    if failed:
        body_lines.append("Failed cases: " + ", ".join(failed[:PRACTICAL_SUMMARY_LIMIT]))
    chunks = _chunks_from_text("\n".join(body_lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "trx",
            "test_count": len(results),
            "failure_count": len(failed),
        },
    )


def _extract_tap_report(path: Path) -> ExtractionResult:
    result_lines: list[str] = []
    failure_count = 0
    for raw_line in _read_text_limited(path).splitlines():
        line = raw_line.strip()
        if line.startswith("ok ") or line.startswith("not ok "):
            result_lines.append(line)
            if line.startswith("not ok "):
                failure_count += 1
    chunks = _chunks_from_text("\n".join(result_lines[:PRACTICAL_SUMMARY_LIMIT]), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "tap",
            "test_count": len(result_lines),
            "failure_count": failure_count,
        },
    )


def _extract_coverage_xml_report(path: Path, *, root: ElementTree.Element | None = None) -> ExtractionResult:
    root = root if root is not None else ElementTree.fromstring(_read_text_limited(path))
    covered = _int_or_none(root.get("lines-covered"))
    valid = _int_or_none(root.get("lines-valid"))
    if covered is not None and valid:
        percent = round((covered / valid) * 100, 2)
    else:
        line_rate = root.get("line-rate")
        percent = round(float(line_rate) * 100, 2) if line_rate is not None else 0.0
    body = f"Line coverage: {percent}%"
    if covered is not None and valid is not None:
        body = f"Lines covered: {covered}/{valid}\n{body}"
    chunks = _chunks_from_text(body, path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "coverage_xml",
            "covered_line_count": covered,
            "line_count": valid,
            "line_coverage_percent": percent,
        },
    )


def _extract_lcov_report(path: Path) -> ExtractionResult:
    covered = 0
    total = 0
    files: set[str] = set()
    for raw_line in _read_text_limited(path).splitlines():
        line = raw_line.strip()
        if line.startswith("SF:"):
            files.add(line[3:])
        elif line.startswith("DA:"):
            total += 1
            fields = line[3:].split(",", 1)
            if len(fields) == 2 and (_int_or_none(fields[1]) or 0) > 0:
                covered += 1
    percent = round((covered / total) * 100, 2) if total else 0.0
    body = f"Files: {len(files)}\nLines covered: {covered}/{total}\nLine coverage: {percent}%"
    chunks = _chunks_from_text(body, path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "report",
            "report_format": "lcov",
            "covered_line_count": covered,
            "line_count": total,
            "line_coverage_percent": percent,
        },
    )


def _extract_bounded_report_text(path: Path) -> ExtractionResult:
    text = _read_text_limited(path)
    lines = [line.strip() for line in text.splitlines() if line.strip()][:PRACTICAL_SUMMARY_LIMIT]
    chunks = _chunks_from_text("\n".join(lines), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": "report", "report_format": path.suffix.lower().lstrip(".") or "unknown"},
    )


def _extract_database(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    if ext not in {".db", ".sqlite", ".sqlite3"}:
        return ExtractionResult(
            status="metadata_only",
            metadata={
                "extractor": "database",
                "database_format": ext.lstrip(".") or "unknown",
                "tool_candidates": _domain_tool_candidates("database", ext),
            },
        )
    uri_path = quote(str(_extractor_path(path)).replace("\\", "/"), safe="/:")
    conn = sqlite3.connect(f"file:{uri_path}?mode=ro", uri=True)
    try:
        rows = conn.execute(
            "SELECT name, type FROM sqlite_master WHERE type IN ('table', 'view') AND name NOT LIKE 'sqlite_%' ORDER BY type, name"
        ).fetchall()
        tables: list[dict[str, Any]] = []
        for name, kind in rows[:PRACTICAL_SUMMARY_LIMIT]:
            columns = [str(row[1]) for row in conn.execute(f"PRAGMA table_info({_sqlite_string_literal(str(name))})")]
            tables.append({"name": str(name), "type": str(kind), "columns": columns})
    finally:
        conn.close()
    body = "\n".join(f"{table['name']} ({', '.join(table['columns'])})" for table in tables)
    chunks = _chunks_from_text(body, path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            "extractor": "database",
            "database_format": "sqlite",
            "table_count": len(rows),
            "tables": tables,
            "truncated": len(rows) > PRACTICAL_SUMMARY_LIMIT,
        },
    )


def _extract_domain_metadata(path: Path, file_kind: str) -> ExtractionResult:
    ext = path.suffix.lower()
    return ExtractionResult(
        status="metadata_only",
        metadata={
            "extractor": file_kind,
            "source_extension": ext,
            "tool_candidates": _domain_tool_candidates(file_kind, ext),
            "metadata_first": True,
        },
    )


def _extract_sensitive_metadata(path: Path) -> ExtractionResult:
    return ExtractionResult(
        status="metadata_only",
        metadata={
            "extractor": "sensitive_metadata",
            "source_extension": path.suffix.lower(),
            "sensitive": True,
            "metadata_first": True,
        },
    )


def _read_text_limited(path: Path, limit: int = PRACTICAL_TEXT_LIMIT_BYTES) -> str:
    data = _read_limited_file(path, limit)
    return decode_text_bytes(data)


def _clean_subtitle_line(line: str) -> str:
    line = re.sub(r"\{\\.*?\}", "", line)
    line = re.sub(r"<[^>]+>", "", line)
    return unescape(line).replace("\\N", "\n").strip()


def _ass_dialogue_text(line: str) -> str:
    _, _, payload = line.partition(":")
    fields = payload.split(",", 9)
    return fields[-1] if fields else payload


def _parse_email_message(data: bytes) -> Any:
    return BytesParser(policy=email_policy.default).parsebytes(data)


def _mail_messages(path: Path) -> list[Any]:
    messages: list[Any] = []
    try:
        box = mailbox.mbox(str(path), create=False)
        try:
            messages = list(box)
        finally:
            box.close()
    except Exception:
        messages = []
    if messages:
        return messages
    return [_parse_email_message(raw) for raw in _split_mbox_messages(read_text_with_bom(path))]


def _split_mbox_messages(text: str) -> list[bytes]:
    parts: list[list[str]] = []
    current: list[str] = []
    for line in text.splitlines():
        if line.startswith("From ") and current:
            parts.append(current)
            current = []
            continue
        if not line.startswith("From "):
            current.append(line)
    if current:
        parts.append(current)
    return ["\n".join(part).encode("utf-8", errors="replace") for part in parts if part]


def _mail_attachment_count(message: Any) -> int:
    count = 0
    for part in message.walk() if hasattr(message, "walk") else ():
        disposition = str(part.get_content_disposition() or "").lower()
        if disposition == "attachment" or part.get_filename():
            count += 1
    return count


def _mail_plain_text(message: Any) -> str:
    if hasattr(message, "get_body"):
        body = message.get_body(preferencelist=("plain",))
        if body is not None:
            try:
                return str(body.get_content()).strip()
            except Exception:
                pass
    parts: list[str] = []
    walk = message.walk() if hasattr(message, "walk") else [message]
    for part in walk:
        if part.is_multipart():
            continue
        disposition = str(part.get_content_disposition() or "").lower()
        if disposition == "attachment" or part.get_filename():
            continue
        content_type = str(part.get_content_type() or "")
        if content_type != "text/plain":
            continue
        payload = part.get_payload(decode=True)
        if isinstance(payload, bytes):
            parts.append(payload.decode(part.get_content_charset() or "utf-8", errors="replace"))
        else:
            parts.append(str(part.get_payload() or ""))
    return "\n".join(part.strip() for part in parts if part.strip())


def _ical_blocks(text: str, block_name: str) -> list[dict[str, list[str]]]:
    target_begin = f"BEGIN:{block_name}"
    target_end = f"END:{block_name}"
    blocks: list[dict[str, list[str]]] = []
    current: dict[str, list[str]] | None = None
    for line in _unfold_ical_lines(text):
        upper = line.upper()
        if upper == target_begin:
            current = {}
            continue
        if upper == target_end and current is not None:
            blocks.append(current)
            current = None
            continue
        if current is None or ":" not in line:
            continue
        key, _, value = line.partition(":")
        prop = key.split(";", 1)[0].upper()
        current.setdefault(prop, []).append(_decode_ical_text(value))
    return blocks


def _unfold_ical_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw_line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        if raw_line.startswith((" ", "\t")) and lines:
            lines[-1] += raw_line[1:]
        else:
            lines.append(raw_line)
    return [line.strip() for line in lines if line.strip()]


def _decode_ical_text(value: str) -> str:
    return value.replace("\\n", "\n").replace("\\N", "\n").replace("\\,", ",").replace("\\;", ";").strip()


def _first_property_value(block: dict[str, list[str]], key: str) -> str | None:
    values = block.get(key)
    if not values:
        return None
    return values[0].strip() or None


def _sqlite_string_literal(value: str) -> str:
    return "'" + value.replace("'", "''") + "'"


def _domain_tool_candidates(file_kind: str, ext: str) -> list[str]:
    if file_kind == "database":
        if ext == ".duckdb":
            return ["duckdb"]
        if ext in {".mdb", ".accdb"}:
            return ["mdbtools"]
        if ext == ".dbf":
            return ["pyarrow", "ogrinfo"]
    if file_kind == "geospatial":
        return ["ogrinfo", "gdalinfo"]
    if file_kind == "cad":
        if ext in {".ifc", ".ifczip"}:
            return ["ifcopenshell"]
        if ext in {".fbx", ".dae", ".obj", ".stl", ".gltf", ".glb"}:
            return ["assimp", "blender"]
        return ["assimp", "blender", "exiftool"]
    if file_kind == "scientific":
        return ["pyarrow", "duckdb"]
    return []


def _extract_document(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".epub":
        return _extract_epub(path, policy)
    if ext == ".fb2":
        return _extract_fb2(path, policy)
    if ext in CALIBRE_PUBLICATION_EXTENSIONS:
        return _extract_calibre_publication(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext in LEGACY_WORD_EXTENSIONS:
        return _extract_legacy_word_document(path)
    if ext in CONVERTED_WORD_EXTENSIONS:
        return _extract_converted_word_document(path)
    if ext in OPENPYXL_EXTENSIONS:
        return _extract_xlsx(path, extractor=ext.lstrip("."))
    if ext in LEGACY_SPREADSHEET_EXTENSIONS:
        return _extract_legacy_spreadsheet(path, policy)
    if ext in OPENDOCUMENT_SPREADSHEET_EXTENSIONS:
        return _extract_opendocument_spreadsheet(path, policy)
    if ext in PPTX_PACKAGE_EXTENSIONS:
        return _extract_pptx(path, extractor=ext.lstrip("."))
    if ext in LEGACY_PRESENTATION_EXTENSIONS:
        return _extract_legacy_presentation(path)
    if ext in OPENDOCUMENT_PRESENTATION_EXTENSIONS:
        return _extract_opendocument_presentation(path)
    return ExtractionResult(status="metadata_only", metadata={"extractor": "document", "extension": ext})


def _extract_pdf(path: Path) -> ExtractionResult:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": "pdf"}, message="pypdf not installed")
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        if exc.__class__.__name__ == "DependencyError":
            message = str(exc) or "pypdf missing required dependency"
            metadata = {"extractor": "pdf"}
            if "cryptography" in message.lower():
                metadata["dependency"] = "cryptography"
            return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=message)
        raise
    page_count = len(reader.pages)
    page_texts = _pdf_page_texts(reader)
    text = "\n".join(text for _page_number, text in page_texts if text).strip()
    chunks = _chunks_from_text(text, path.name)
    ocr_pages = _pdf_pages_requiring_ocr(page_texts)
    if chunks and not ocr_pages:
        return ExtractionResult(
            status="indexed",
            chunks=chunks,
            metadata={
                "extractor": "pdf",
                "page_count": page_count,
                "ocr": {
                    "status": "skipped_embedded_text",
                    "page_count": page_count,
                    "pages_attempted": 0,
                    "cache_hits": 0,
                    "cache_misses": 0,
                },
            },
        )
    if page_count <= 0:
        return ExtractionResult(
            status="metadata_only",
            metadata={
                "extractor": "pdf",
                "page_count": 0,
                "ocr": {"status": "completed", "page_count": 0, "pages_attempted": 0, "cache_hits": 0, "cache_misses": 0},
            },
        )
    if not ocr_pages:
        ocr_pages = list(range(1, page_count + 1))
    ocr = _ocr_pdf_pages(path, page_start=ocr_pages[0], page_end=ocr_pages[-1], page_count=page_count, page_numbers=ocr_pages)
    metadata = {"extractor": "pdf", "page_count": page_count, "ocr": ocr.metadata}
    if ocr.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", chunks=chunks, metadata=metadata, message=ocr.message)
    if ocr.status == "failed":
        return ExtractionResult(status="failed", chunks=chunks, metadata=metadata, message=ocr.message)
    ocr_chunks = _offset_chunks(
        _chunks_from_text(ocr.text, path.name, modality="ocr"),
        PDF_OCR_CHUNK_INDEX_BASE,
        locator_prefix=_pdf_page_locator(ocr_pages),
    )
    all_chunks = tuple(chunks) + ocr_chunks
    return ExtractionResult(status="indexed" if all_chunks else "metadata_only", chunks=all_chunks, metadata=metadata, message=ocr.message)


def plan_staged_pdf_extraction(path: str | Path, _policy: CorpusPolicy | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    try:
        from pypdf import PdfReader
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": "pdf"}, message="pypdf not installed")
    try:
        reader = PdfReader(str(file_path))
    except Exception as exc:
        if exc.__class__.__name__ == "DependencyError":
            message = str(exc) or "pypdf missing required dependency"
            metadata = {"extractor": "pdf"}
            if "cryptography" in message.lower():
                metadata["dependency"] = "cryptography"
            return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=message)
        raise
    page_count = len(reader.pages)
    page_texts = _pdf_page_texts(reader)
    text = "\n".join(text for _page_number, text in page_texts if text).strip()
    chunks = _chunks_from_text(text, file_path.name)
    ocr_pages = _pdf_pages_requiring_ocr(page_texts)
    if chunks and not ocr_pages:
        return ExtractionResult(
            status="indexed",
            chunks=chunks,
            metadata={
                "extractor": "pdf",
                "page_count": page_count,
                "ocr": {
                    "status": "skipped_embedded_text",
                    "page_count": page_count,
                    "pages_attempted": 0,
                    "cache_hits": 0,
                    "cache_misses": 0,
                },
            },
        )
    if page_count <= 0:
        return ExtractionResult(
            status="metadata_only",
            metadata={
                "extractor": "pdf",
                "page_count": 0,
                "ocr": {"status": "completed", "page_count": 0, "pages_attempted": 0, "cache_hits": 0, "cache_misses": 0},
            },
        )
    if not ocr_pages:
        ocr_pages = list(range(1, page_count + 1))
    page_batch = ocr_pages[:OCR_PDF_PAGE_BATCH_SIZE]
    remaining_pages = ocr_pages[OCR_PDF_PAGE_BATCH_SIZE:]
    staged_job = {
        "job_type": "corpus_extract_pdf_ocr_pages",
        "payload": {
            "pages": page_batch,
            "page_count": page_count,
            "page_batch_size": OCR_PDF_PAGE_BATCH_SIZE,
            "chunks_seen": len(chunks),
            "embedded_chunk_count": len(chunks),
        },
    }
    if remaining_pages:
        staged_job["payload"]["remaining_pages"] = remaining_pages
    batch_count = math.ceil(len(ocr_pages) / OCR_PDF_PAGE_BATCH_SIZE)
    metadata = {
        "extractor": "pdf",
        "page_count": page_count,
        "ocr": {
            "engine": "tesseract",
            "renderer": "pdftoppm",
            "status": "planned",
            "page_count": page_count,
            "pages_attempted": 0,
            "pages_planned": len(ocr_pages),
            "pages_with_embedded_text": page_count - len(ocr_pages),
            "cache_hits": 0,
            "cache_misses": 0,
            "page_batch_size": OCR_PDF_PAGE_BATCH_SIZE,
        },
        "staged_extraction": {
            "status": "planned",
            "content_extracted": bool(chunks),
            "pending_job_count": batch_count,
            "next_job_type": staged_job["job_type"],
            "unit": "pdf_pages",
            "page_count": page_count,
            "page_batch_size": OCR_PDF_PAGE_BATCH_SIZE,
        },
        "staged_jobs": [staged_job],
    }
    return ExtractionResult(status="staged", chunks=chunks, metadata=metadata)


def extract_pdf_ocr_pages(path: str | Path, payload: dict[str, Any] | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    payload = payload or {}
    page_count = _optional_int(payload.get("page_count"))
    if page_count is None:
        try:
            from pypdf import PdfReader
        except ImportError:
            return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": "pdf"}, message="pypdf not installed")
        page_count = len(PdfReader(str(file_path)).pages)
    if page_count <= 0:
        metadata = {
            "extractor": "pdf_ocr_pages",
            "page_count": 0,
            "ocr": {"status": "completed", "page_count": 0, "pages_attempted": 0, "cache_hits": 0, "cache_misses": 0},
            "staged_extraction": {
                "status": "piece_completed",
                "complete": True,
                "unit": "pdf_pages",
                "page_count": 0,
                "chunks_written": 0,
                "chunks_seen": max(0, _optional_int(payload.get("chunks_seen")) or 0),
            },
        }
        return ExtractionResult(status="metadata_only", metadata=metadata)
    explicit_pages = _normalised_page_numbers(payload.get("pages"), page_count)
    page_start = min(explicit_pages) if explicit_pages else max(1, _optional_int(payload.get("page_start")) or 1)
    page_batch_size = max(1, _optional_int(payload.get("page_batch_size")) or OCR_PDF_PAGE_BATCH_SIZE)
    page_end = max(explicit_pages) if explicit_pages else min(page_count, _optional_int(payload.get("page_end")) or (page_start + page_batch_size - 1))
    chunks_seen = max(0, _optional_int(payload.get("chunks_seen")) or 0)
    ocr = _ocr_pdf_pages(file_path, page_start=page_start, page_end=page_end, page_count=page_count, page_numbers=explicit_pages or None)
    metadata = {
        "extractor": "pdf_ocr_pages",
        "page_count": page_count,
        "ocr": ocr.metadata,
    }
    if ocr.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=ocr.message)
    if ocr.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=ocr.message)
    page_label = _pdf_page_locator(explicit_pages or list(range(page_start, page_end + 1)))
    local_chunks = _chunks_from_text(ocr.text, f"{file_path.name} {page_label}", modality="ocr")
    chunks = _offset_chunks(
        local_chunks,
        PDF_OCR_CHUNK_INDEX_BASE + ((page_start - 1) * MEDIA_SEGMENT_CHUNK_INDEX_STRIDE),
        locator_prefix=page_label,
    )
    total_chunks_seen = chunks_seen + len(chunks)
    next_job: dict[str, Any] | None = None
    remaining_pages = _normalised_page_numbers(payload.get("remaining_pages"), page_count)
    if remaining_pages:
        next_pages = remaining_pages[:page_batch_size]
        next_remaining = remaining_pages[page_batch_size:]
        next_job = {
            "job_type": "corpus_extract_pdf_ocr_pages",
            "payload": {
                "pages": next_pages,
                "page_count": page_count,
                "page_batch_size": page_batch_size,
                "chunks_seen": total_chunks_seen,
                "embedded_chunk_count": _optional_int(payload.get("embedded_chunk_count")) or 0,
            },
        }
        if next_remaining:
            next_job["payload"]["remaining_pages"] = next_remaining
    elif not explicit_pages and page_end < page_count:
        next_start = page_end + 1
        next_end = min(page_count, next_start + page_batch_size - 1)
        next_job = {
            "job_type": "corpus_extract_pdf_ocr_pages",
            "payload": {
                "page_start": next_start,
                "page_end": next_end,
                "page_count": page_count,
                "page_batch_size": page_batch_size,
                "chunks_seen": total_chunks_seen,
            },
        }
    metadata["staged_extraction"] = {
        "status": "piece_completed",
        "complete": next_job is None,
        "unit": "pdf_pages",
        "page_start": page_start,
        "page_end": page_end,
        "pages": explicit_pages or list(range(page_start, page_end + 1)),
        "page_count": page_count,
        "chunks_written": len(chunks),
        "chunks_seen": total_chunks_seen,
    }
    if next_job is not None:
        metadata["staged_extraction"]["next_job"] = next_job
        return ExtractionResult(status="staged", chunks=chunks, metadata=metadata, message=ocr.message)
    return ExtractionResult(
        status="indexed" if total_chunks_seen > 0 else "metadata_only",
        chunks=chunks,
        metadata=metadata,
        message=ocr.message,
    )


def _pdf_page_texts(reader: Any) -> list[tuple[int, str]]:
    return [(index, (page.extract_text() or "").strip()) for index, page in enumerate(reader.pages, start=1)]


def _pdf_pages_requiring_ocr(page_texts: Iterable[tuple[int, str]]) -> list[int]:
    return [page_number for page_number, text in page_texts if not text.strip()]


def _normalised_page_numbers(value: Any, page_count: int) -> list[int]:
    if not isinstance(value, list):
        return []
    pages: list[int] = []
    for item in value:
        page_number = _optional_int(item)
        if page_number is not None and 1 <= page_number <= page_count and page_number not in pages:
            pages.append(page_number)
    return pages


def _pdf_page_locator(pages: list[int]) -> str:
    if not pages:
        return "page:unknown"
    if len(pages) == 1:
        return f"page:{pages[0]}"
    if pages == list(range(pages[0], pages[-1] + 1)):
        return f"page:{pages[0]}-{pages[-1]}"
    preview = ",".join(str(page) for page in pages[:8])
    if len(pages) > 8:
        preview += ",..."
    return f"page:{preview}"


def _extract_epub(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    metadata = _publication_metadata(path, "epub")
    try:
        with ZipFile(path) as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
            total_bytes = 0
            parts: list[str] = []
            content_file_count = 0
            for index, info in enumerate(infos):
                member_path = _safe_container_member_name(info.filename)
                if member_path is None:
                    return ExtractionResult(status="failed", metadata=metadata, message=f"unsafe EPUB member: {info.filename}")
                if info.flag_bits & 0x1:
                    return ExtractionResult(status="failed", metadata=metadata, message=f"encrypted EPUB member is not supported: {member_path}")
                size = int(info.file_size or 0)
                cap_message = _container_cap_message(policy, member_count=index + 1, member_size=size, total_bytes=total_bytes + size)
                if cap_message:
                    return ExtractionResult(status="metadata_only", metadata={**metadata, "warnings": [cap_message]}, message=cap_message)
                data = archive.read(info)
                total_bytes += len(data)
                lower_member = member_path.lower()
                if lower_member.endswith(".opf"):
                    metadata.update(_epub_package_metadata(data))
                elif lower_member.endswith((".xhtml", ".html", ".htm")):
                    text = _text_from_markup_bytes(data)
                    if text:
                        parts.append(text)
                        content_file_count += 1
    except BadZipFile as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=f"EPUB parse failed: {exc}")
    except ValueError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc))

    metadata["content_file_count"] = content_file_count
    chunks = _chunks_from_text("\n\n".join(parts), path.name)
    return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata=metadata)


def _extract_fb2(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    metadata = _publication_metadata(path, "fb2")
    size = path.stat().st_size
    if size > policy.container_max_member_bytes:
        message = "FB2 file exceeds publication extraction size limit"
        return ExtractionResult(status="metadata_only", metadata={**metadata, "warnings": [message]}, message=message)
    try:
        root = ElementTree.fromstring(path.read_bytes())
    except ElementTree.ParseError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=f"FB2 parse failed: {exc}")

    title = _first_local_text(root, "book-title")
    author = _fb2_author(root)
    if title:
        metadata["publication_title"] = title
    if author:
        metadata["publication_author"] = author
    body = _first_local_element(root, "body")
    paragraphs = [_normalized_text(paragraph.itertext()) for paragraph in (body.iter() if body is not None else ()) if _xml_local_name(paragraph.tag) == "p"]
    paragraphs = [paragraph for paragraph in paragraphs if paragraph]
    metadata["paragraph_count"] = len(paragraphs)
    metadata["content_file_count"] = 1 if paragraphs else 0
    chunks = _chunks_from_text("\n\n".join(paragraphs), path.name)
    return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata=metadata)


def _extract_calibre_publication(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    metadata = _publication_metadata(path, ext.lstrip("."))
    command = shutil.which("ebook-convert")
    if command is None:
        return ExtractionResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "attempted": ["ebook-convert"]},
            message="MOBI/AZW/LIT extraction requires Calibre ebook-convert.",
        )
    with tempfile.TemporaryDirectory(prefix="flux-kb-ebook-") as temp_dir:
        output_path = Path(temp_dir) / f"{path.stem}.txt"
        try:
            result = run_no_window(
                [command, str(path), str(output_path)],
                text=True,
                capture_output=True,
                timeout=180,
                check=False,
            )
        except Exception as exc:  # pragma: no cover - environment-specific
            return ExtractionResult(status="failed", metadata={**metadata, "attempted": ["ebook-convert"]}, message=str(exc))
        if result.returncode != 0:
            message = result.stderr.strip() if isinstance(result.stderr, str) else str(result.stderr or "ebook-convert failed")
            return ExtractionResult(status="failed", metadata={**metadata, "attempted": ["ebook-convert"]}, message=message)
        text = read_text_with_bom(output_path).strip() if output_path.exists() else ""
    chunks = _chunks_from_text(text, path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={
            **metadata,
            "extractor": "ebook_convert",
            "source_extension": ext,
            "converted_extension": ".txt",
        },
    )


def _publication_metadata(path: Path, publication_format: str) -> dict[str, Any]:
    return {
        "extractor": "publication",
        "publication_type": "ebook",
        "publication_format": publication_format,
        "source_extension": path.suffix.lower(),
    }


def _epub_package_metadata(data: bytes) -> dict[str, Any]:
    try:
        root = ElementTree.fromstring(data)
    except ElementTree.ParseError:
        return {}
    metadata: dict[str, Any] = {}
    title = _first_local_text(root, "title")
    creator = _first_local_text(root, "creator")
    if title:
        metadata["publication_title"] = title
    if creator:
        metadata["publication_author"] = creator
    return metadata


def _fb2_author(root: ElementTree.Element) -> str | None:
    for element in root.iter():
        if _xml_local_name(element.tag) != "author":
            continue
        parts = [
            _normalized_text(child.itertext())
            for child in element
            if _xml_local_name(child.tag) in {"first-name", "middle-name", "last-name", "nickname"}
        ]
        author = " ".join(part for part in parts if part)
        if author:
            return author
    return None


def _text_from_markup_bytes(data: bytes) -> str:
    raw = data.decode("utf-8", errors="replace")
    try:
        root = ElementTree.fromstring(raw)
    except ElementTree.ParseError:
        stripped = re.sub(r"<[^>]+>", " ", raw)
        return _normalized_text([unescape(stripped)])
    return _normalized_text(root.itertext())


def _first_local_element(root: ElementTree.Element, local_name: str) -> ElementTree.Element | None:
    for element in root.iter():
        if _xml_local_name(element.tag) == local_name:
            return element
    return None


def _first_local_text(root: ElementTree.Element, local_name: str) -> str | None:
    element = _first_local_element(root, local_name)
    if element is None:
        return None
    text = _normalized_text(element.itertext())
    return text or None


def _normalized_text(parts: Iterable[str]) -> str:
    return re.sub(r"\s+", " ", " ".join(str(part) for part in parts)).strip()


def _xml_local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _extract_docx(path: Path, *, extractor: str = "docx") -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": extractor}, message="python-docx not installed")
    try:
        document = Document(str(path))
    except Exception as exc:
        missing_package_part = _missing_archive_member_from_key_error(exc) if isinstance(exc, KeyError) else None
        if missing_package_part is not None:
            return _extract_docx_package_xml(
                path,
                extractor=extractor,
                missing_package_part=missing_package_part,
                warning=str(exc),
            )
        if _is_invalid_package_error(exc):
            return ExtractionResult(
                status="blocked_missing_dependency",
                metadata={"extractor": extractor, "reason": "invalid_package"},
                message=str(exc),
            )
        raise
    text = "\n".join(paragraph.text for paragraph in document.paragraphs).strip()
    chunks = _chunks_from_text(text, path.name)
    return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata={"extractor": extractor})


def _missing_archive_member_from_key_error(exc: KeyError) -> str | None:
    match = re.search(r"There is no item named '([^']+)' in the archive", str(exc))
    return match.group(1) if match else None


def _extract_docx_package_xml(
    path: Path,
    *,
    extractor: str,
    missing_package_part: str,
    warning: str,
) -> ExtractionResult:
    metadata = {
        "extractor": extractor,
        "fallback": "package_xml",
        "missing_package_part": missing_package_part,
        "warnings": [warning],
    }
    try:
        with ZipFile(path) as archive:
            document_xml = archive.read("word/document.xml")
    except (BadZipFile, KeyError) as exc:
        return ExtractionResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "reason": "invalid_package"},
            message=str(exc),
        )
    try:
        root = ElementTree.fromstring(document_xml)
    except ElementTree.ParseError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=f"DOCX XML parse failed: {exc}")

    paragraphs = [
        paragraph_text
        for paragraph in root.iter()
        if _xml_local_name(paragraph.tag) == "p"
        for paragraph_text in [_docx_paragraph_text(paragraph)]
        if paragraph_text
    ]
    chunks = _chunks_from_text("\n".join(paragraphs), path.name)
    return ExtractionResult(status="indexed" if chunks else "metadata_only", chunks=chunks, metadata=metadata)


def _docx_paragraph_text(paragraph: ElementTree.Element) -> str:
    parts: list[str] = []
    for element in paragraph.iter():
        local_name = _xml_local_name(element.tag)
        if local_name == "t" and element.text:
            parts.append(element.text)
        elif local_name == "tab":
            parts.append("\t")
        elif local_name in {"br", "cr"}:
            parts.append("\n")
    return "".join(parts).strip()


def _is_invalid_package_error(exc: Exception) -> bool:
    message = str(exc).lower()
    return exc.__class__.__name__ in {"PackageNotFoundError", "BadZipFile"} or "package not found" in message


def _invalid_package_result(exc: Exception, *, extractor: str) -> ExtractionResult:
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": extractor, "reason": "invalid_package"},
        message=str(exc),
    )


def _extract_pptx(path: Path, *, extractor: str = "pptx") -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": extractor}, message="python-pptx not installed")
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    chunks = _chunks_from_text("\n".join(parts), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": extractor, "slide_count": len(presentation.slides)},
    )


def _extract_xlsx(path: Path, *, extractor: str = "xlsx") -> ExtractionResult:
    try:
        import openpyxl
    except ImportError:
        return ExtractionResult(status="blocked_missing_dependency", metadata={"extractor": extractor}, message="openpyxl not installed")
    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        if _is_invalid_package_error(exc):
            return _invalid_package_result(exc, extractor=extractor)
        raise
    parts: list[str] = []
    for worksheet in workbook.worksheets[:10]:
        parts.append(f"Sheet: {worksheet.title}")
        for row in worksheet.iter_rows(max_row=200, max_col=30, values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(" | ".join(values))
    chunks = _chunks_from_text("\n".join(parts), path.name)
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata={"extractor": extractor, "sheet_count": len(workbook.worksheets)},
    )


def _extract_converted_word_document(path: Path) -> ExtractionResult:
    result = _extract_via_libreoffice_conversion(
        path,
        target_format="docx",
        target_suffix=".docx",
        read_converted=lambda converted_path: _extract_docx(converted_path),
    )
    if result is not None:
        return result
    ext = path.suffix.lower()
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": "word_document", "extension": ext, "attempted": ["libreoffice"]},
        message="Word-like document extraction requires LibreOffice for this format.",
    )


def _extract_legacy_spreadsheet(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    ext = path.suffix.lower()
    result = _extract_via_libreoffice_conversion(
        path,
        target_format="xlsx",
        target_suffix=".xlsx",
        read_converted=lambda converted_path: _extract_converted_spreadsheet_workbook(converted_path, policy),
    )
    if result is not None and result.status != "blocked_missing_dependency":
        return result
    text = _extract_with_excel_com(path)
    if text:
        chunks = _chunks_from_text(text, path.name)
        return ExtractionResult(
            status="indexed" if chunks else "metadata_only",
            chunks=chunks,
            metadata={"extractor": "excel_com", "source_extension": ext},
        )
    if result is not None:
        return result
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": "legacy_spreadsheet", "extension": ext, "attempted": ["libreoffice", "excel_com"]},
        message="Legacy spreadsheet extraction requires LibreOffice or Windows Excel COM.",
    )


def _extract_opendocument_spreadsheet(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    ext = path.suffix.lower()
    result = _extract_via_libreoffice_conversion(
        path,
        target_format="xlsx",
        target_suffix=".xlsx",
        read_converted=lambda converted_path: _extract_converted_spreadsheet_workbook(converted_path, policy),
    )
    if result is not None:
        return result
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": "opendocument_spreadsheet", "extension": ext, "attempted": ["libreoffice"]},
        message="OpenDocument spreadsheet extraction requires LibreOffice.",
    )


def _extract_converted_spreadsheet_workbook(path: Path, policy: CorpusPolicy) -> ExtractionResult:
    if path.stat().st_size > policy.max_inline_bytes:
        return _extract_sample_first_workbook(path, extractor="sample_first_workbook")
    return _extract_xlsx(path)


def _extract_legacy_presentation(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    result = _extract_via_libreoffice_conversion(
        path,
        target_format="pptx",
        target_suffix=".pptx",
        read_converted=lambda converted_path: _extract_pptx(converted_path),
    )
    if result is not None and result.status != "blocked_missing_dependency":
        return result
    text = _extract_with_powerpoint_com(path)
    if text:
        chunks = _chunks_from_text(text, path.name)
        return ExtractionResult(
            status="indexed" if chunks else "metadata_only",
            chunks=chunks,
            metadata={"extractor": "powerpoint_com", "source_extension": ext},
        )
    if result is not None:
        return result
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": "legacy_presentation", "extension": ext, "attempted": ["libreoffice", "powerpoint_com"]},
        message="Legacy presentation extraction requires LibreOffice or Windows PowerPoint COM.",
    )


def _extract_opendocument_presentation(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    result = _extract_via_libreoffice_conversion(
        path,
        target_format="pptx",
        target_suffix=".pptx",
        read_converted=lambda converted_path: _extract_pptx(converted_path),
    )
    if result is not None:
        return result
    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={"extractor": "opendocument_presentation", "extension": ext, "attempted": ["libreoffice"]},
        message="OpenDocument presentation extraction requires LibreOffice.",
    )


def _extract_via_libreoffice_conversion(
    path: Path,
    *,
    target_format: str,
    target_suffix: str,
    read_converted: Callable[[Path], ExtractionResult],
) -> ExtractionResult | None:
    with tempfile.TemporaryDirectory(prefix="flux-kb-lo-") as temp_dir:
        converted_path = _convert_with_libreoffice(
            path,
            target_format=target_format,
            target_suffix=target_suffix,
            output_dir=Path(temp_dir),
        )
        if converted_path is None:
            return None
        result = read_converted(converted_path)
        metadata = {
            **result.metadata,
            "extractor": "libreoffice",
            "source_extension": path.suffix.lower(),
            "converted_extension": target_suffix,
        }
        return ExtractionResult(
            status=result.status,
            chunks=result.chunks,
            metadata=metadata,
            message=result.message,
        )


def _convert_with_libreoffice(
    path: Path,
    *,
    target_format: str,
    target_suffix: str,
    output_dir: Path,
) -> Path | None:
    command = shutil.which("soffice") or shutil.which("libreoffice")
    if command is None:
        return None
    try:
        result = run_no_window(
            [
                command,
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(output_dir),
                str(path),
            ],
            text=True,
            capture_output=True,
            timeout=120,
            check=False,
        )
        if result.returncode != 0:
            return None
        output_path = output_dir / f"{path.stem}{target_suffix}"
        candidates = [output_path] if output_path.exists() else sorted(output_dir.glob(f"*{target_suffix}"))
        return candidates[0] if candidates else None
    except Exception:
        return None


def _extract_legacy_word_document(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    metadata = {"extractor": "legacy_document", "extension": ext}
    attempts: list[str] = []

    for extractor_name, extract in (
        ("libreoffice", _extract_with_libreoffice),
        ("antiword", _extract_with_antiword),
        ("catdoc", _extract_with_catdoc),
        ("wvText", _extract_with_wvtext),
        ("word_com", _extract_with_word_com),
    ):
        attempts.append(extractor_name)
        text = extract(path)
        if not text:
            continue
        chunks = _chunks_from_text(text, path.name)
        return ExtractionResult(
            status="indexed" if chunks else "metadata_only",
            chunks=chunks,
            metadata={"extractor": extractor_name, "extension": ext},
        )

    return ExtractionResult(
        status="blocked_missing_dependency",
        metadata={**metadata, "attempted": attempts},
        message="Legacy Word extraction requires LibreOffice, antiword, catdoc, wvText, or Windows Word COM.",
    )


def _extract_with_libreoffice(path: Path) -> str | None:
    command = shutil.which("soffice") or shutil.which("libreoffice")
    if command is None:
        return None
    try:
        with tempfile.TemporaryDirectory(prefix="flux-kb-lo-") as temp_dir:
            result = run_no_window(
                [
                    command,
                    "--headless",
                    "--convert-to",
                    "txt:Text",
                    "--outdir",
                    temp_dir,
                    str(path),
                ],
                text=True,
                capture_output=True,
                timeout=90,
                check=False,
            )
            if result.returncode != 0:
                return None
            output_path = Path(temp_dir) / f"{path.stem}.txt"
            candidates = [output_path] if output_path.exists() else sorted(Path(temp_dir).glob("*.txt"))
            for candidate in candidates:
                text = read_text_with_bom(candidate).strip()
                if text:
                    return text
    except Exception:
        return None
    return None


def _extract_with_antiword(path: Path) -> str | None:
    return _extract_with_stdout_tool(path, "antiword")


def _extract_with_catdoc(path: Path) -> str | None:
    return _extract_with_stdout_tool(path, "catdoc")


def _extract_with_stdout_tool(path: Path, command_name: str) -> str | None:
    command = shutil.which(command_name)
    if command is None:
        return None
    try:
        result = run_no_window(
            [command, str(path)],
            text=True,
            capture_output=True,
            timeout=60,
            check=False,
        )
    except Exception:
        return None
    if result.returncode != 0:
        return None
    return result.stdout.strip() or None


def _extract_with_wvtext(path: Path) -> str | None:
    command = shutil.which("wvText")
    if command is None:
        return None
    try:
        with tempfile.TemporaryDirectory(prefix="flux-kb-wv-") as temp_dir:
            output_path = Path(temp_dir) / f"{path.stem}.txt"
            result = run_no_window(
                [command, str(path), str(output_path)],
                text=True,
                capture_output=True,
                timeout=60,
                check=False,
            )
            if result.returncode != 0 or not output_path.exists():
                return None
            return read_text_with_bom(output_path).strip() or None
    except Exception:
        return None


def _extract_with_word_com(path: Path) -> str | None:
    if os.name != "nt":
        return None
    try:
        import pythoncom
        import win32com.client
    except Exception:
        return None

    word = None
    document = None
    initialized = False
    try:
        pythoncom.CoInitialize()
        initialized = True
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        document = word.Documents.Open(
            FileName=str(path),
            ConfirmConversions=False,
            ReadOnly=True,
            AddToRecentFiles=False,
            Visible=False,
            OpenAndRepair=True,
        )
        return str(document.Content.Text or "").strip() or None
    except Exception:
        return None
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        if initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _extract_with_excel_com(path: Path) -> str | None:
    if os.name != "nt":
        return None
    try:
        import pythoncom
        import win32com.client
    except Exception:
        return None

    excel = None
    workbook = None
    initialized = False
    try:
        pythoncom.CoInitialize()
        initialized = True
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        workbook = excel.Workbooks.Open(
            Filename=str(path),
            ReadOnly=True,
            UpdateLinks=0,
            AddToMru=False,
        )
        parts: list[str] = []
        sheet_count = min(int(workbook.Worksheets.Count), 10)
        for sheet_index in range(1, sheet_count + 1):
            worksheet = workbook.Worksheets(sheet_index)
            parts.append(f"Sheet: {worksheet.Name}")
            for row in _com_table_rows(worksheet.UsedRange.Value, max_rows=200, max_cols=30):
                if row:
                    parts.append(" | ".join(row))
        return "\n".join(parts).strip() or None
    except Exception:
        return None
    finally:
        if workbook is not None:
            try:
                workbook.Close(False)
            except Exception:
                pass
        if excel is not None:
            try:
                excel.Quit()
            except Exception:
                pass
        if initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _extract_with_powerpoint_com(path: Path) -> str | None:
    if os.name != "nt":
        return None
    try:
        import pythoncom
        import win32com.client
    except Exception:
        return None

    powerpoint = None
    presentation = None
    initialized = False
    try:
        pythoncom.CoInitialize()
        initialized = True
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            FileName=str(path),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False,
        )
        parts: list[str] = []
        slide_count = min(int(presentation.Slides.Count), 200)
        for slide_index in range(1, slide_count + 1):
            slide = presentation.Slides(slide_index)
            slide_parts: list[str] = []
            for shape_index in range(1, int(slide.Shapes.Count) + 1):
                shape = slide.Shapes(shape_index)
                if bool(getattr(shape, "HasTextFrame", False)) and bool(shape.TextFrame.HasText):
                    text = str(shape.TextFrame.TextRange.Text or "").strip()
                    if text:
                        slide_parts.append(" ".join(text.split()))
            if slide_parts:
                parts.append(f"Slide {slide_index}: {' '.join(slide_parts)}")
        return "\n".join(parts).strip() or None
    except Exception:
        return None
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        if initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _com_table_rows(value: Any, *, max_rows: int, max_cols: int) -> list[list[str]]:
    if value is None:
        return []
    if not isinstance(value, tuple):
        return [[str(value)]]
    if not value:
        return []
    if not isinstance(value[0], tuple):
        rows = (value,)
    else:
        rows = value
    normalized: list[list[str]] = []
    for row in rows[:max_rows]:
        cells = row if isinstance(row, tuple) else (row,)
        values = [str(cell) for cell in cells[:max_cols] if cell is not None]
        if values:
            normalized.append(values)
    return normalized


ZIP_CONTAINER_EXTENSIONS = {
    ".apk",
    ".cbz",
    ".ear",
    ".egg",
    ".ipa",
    ".jar",
    ".nupkg",
    ".vsix",
    ".war",
    ".whl",
    ".xpi",
    ".zip",
}
TAR_CONTAINER_EXTENSIONS = {".cbt", ".crate", ".gem", ".tar", ".tgz"}
STREAM_CONTAINER_FORMATS = {".bz2": "bzip2", ".gz": "gzip", ".xz": "xz"}
OPTIONAL_CONTAINER_TOOLS = {
    "7z": ("7z", "7zz", "bsdtar", "unar"),
    "cb7": ("7z", "7zz", "bsdtar", "unar"),
    "cbr": ("unrar", "7z", "7zz", "bsdtar", "unar"),
    "rar": ("unrar", "7z", "7zz", "bsdtar", "unar"),
    "cab": ("7z", "7zz", "bsdtar", "unar"),
    "iso": ("7z", "7zz", "bsdtar", "unar"),
    "dmg": ("7z", "7zz", "bsdtar", "unar"),
    "zst": ("zstd", "bsdtar"),
    "lz4": ("lz4", "bsdtar"),
    "ar": ("ar", "bsdtar"),
    "cpio": ("bsdtar",),
    "deb": ("ar", "bsdtar"),
    "rpm": ("rpm2cpio", "bsdtar"),
    "crx": ("7z", "7zz", "bsdtar", "unar"),
}
TEXT_MEMBER_NAMES = {"changelog", "copying", "license", "metadata", "notice", "readme"}
TEXT_MEMBER_EXTENSIONS = TEXT_EXTENSIONS | {".err", ".log", ".out", ".trace"}
EMBEDDED_PARSE_MEMBER_KINDS = {
    "document",
    "diagram",
    "image",
    "audio",
    "video",
    "subtitle",
    "mail",
    "calendar",
    "contact",
    "structured_data",
    "report",
    "database",
    "geospatial",
    "cad",
    "scientific",
    "sensitive_metadata",
}


def _extract_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    depth: int = 0,
    member_prefix: str = "",
) -> ExtractionResult:
    format_name = _container_format(path)
    if format_name == "zip":
        return _with_comic_archive_metadata(
            _extract_zip_container(path, policy, container_kind=container_kind, depth=depth, member_prefix=member_prefix),
            path,
        )
    if format_name == "tar":
        return _with_comic_archive_metadata(
            _extract_tar_container(path, policy, container_kind=container_kind, depth=depth, member_prefix=member_prefix),
            path,
        )
    if format_name in STREAM_CONTAINER_FORMATS.values():
        return _with_comic_archive_metadata(
            _extract_stream_container(path, policy, container_kind=container_kind, format_name=format_name, depth=depth, member_prefix=member_prefix),
            path,
        )
    if format_name in {"zst", "lz4"}:
        stream_result = _extract_tool_stream_container(
            path,
            policy,
            container_kind=container_kind,
            format_name=format_name,
            depth=depth,
            member_prefix=member_prefix,
        )
        if stream_result is not None:
            return _with_comic_archive_metadata(stream_result, path)
    optional_result = _extract_optional_tool_container(
        path,
        policy,
        container_kind=container_kind,
        format_name=format_name,
        depth=depth,
        member_prefix=member_prefix,
    )
    if optional_result is not None:
        return _with_comic_archive_metadata(optional_result, path)
    attempted = list(OPTIONAL_CONTAINER_TOOLS.get(format_name, ()))
    return _with_comic_archive_metadata(
        ExtractionResult(
            status="blocked_missing_dependency",
            metadata=_container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix, attempted=attempted),
            message=f"{format_name} extraction requires a local tool: {', '.join(attempted) or format_name}.",
        ),
        path,
    )


def _container_format(path: Path) -> str:
    name = path.name.lower()
    ext = path.suffix.lower()
    if ext in ZIP_CONTAINER_EXTENSIONS:
        return "zip"
    if name.endswith((".tar.zst", ".tar.lz4")):
        return ext.lstrip(".")
    if ext in TAR_CONTAINER_EXTENSIONS or name.endswith((".tar.gz", ".tar.bz2", ".tar.xz")):
        return "tar"
    if ext in STREAM_CONTAINER_FORMATS:
        return STREAM_CONTAINER_FORMATS[ext]
    return ext.lstrip(".") or "unknown"


def _with_comic_archive_metadata(result: ExtractionResult, path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    if ext not in COMIC_ARCHIVE_EXTENSIONS:
        return result
    return replace(
        result,
        metadata={
            **result.metadata,
            "publication_type": "comic_archive",
            "publication_format": ext.lstrip("."),
        },
    )


def _extract_zip_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    depth: int,
    member_prefix: str,
) -> ExtractionResult:
    format_name = "zip"
    metadata = _container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix)
    try:
        with ZipFile(path) as archive:
            infos = [info for info in archive.infolist() if not info.is_dir()]
            member_payloads: list[dict[str, Any]] = []
            skipped_children: list[ContainerChildAsset] = []
            total_bytes = 0
            for index, info in enumerate(infos):
                member_path = _safe_container_member_name(info.filename)
                if member_path is None:
                    return _failed_container_result(metadata, f"unsafe container member: {info.filename}")
                if info.flag_bits & 0x1:
                    return _failed_container_result(metadata, f"encrypted container member is not supported: {member_path}")
                size = int(info.file_size or 0)
                cap_message = _container_cap_message(policy, member_count=index + 1, member_size=0, total_bytes=total_bytes + size)
                if cap_message:
                    metadata["member_count"] = len(infos)
                    return _metadata_only_container_result(metadata, cap_message)
                if size > policy.container_max_member_bytes:
                    skipped_children.append(
                        _skipped_container_child(
                            _join_container_member_path(member_prefix, member_path),
                            size_bytes=size,
                            reason="member_size_limit",
                            message="member exceeds size limit",
                            container_format=format_name,
                            container_kind=container_kind,
                            member_index=index,
                            compressed_size=int(info.compress_size or 0),
                            member_depth=depth + 1,
                            parent_member_path=member_prefix or None,
                        )
                    )
                    total_bytes += size
                    continue
                total_bytes += size
                data = archive.read(info)
                member_payloads.append(
                    {
                        "member_path": _join_container_member_path(member_prefix, member_path),
                        "data": data,
                        "member_index": index,
                        "compressed_size": int(info.compress_size or 0),
                        "member_depth": depth + 1,
                        "parent_member_path": member_prefix or None,
                    }
                )
            sidecars = _embedded_media_sidecars(member_payloads)
            children: list[ContainerChildAsset] = list(skipped_children)
            for member in member_payloads:
                children.extend(
                    _container_child_assets_from_bytes(
                        member["member_path"],
                        member["data"],
                        policy,
                        container_format=format_name,
                        container_kind=container_kind,
                        member_index=member["member_index"],
                        compressed_size=member["compressed_size"],
                        member_depth=member["member_depth"],
                        parent_member_path=member["parent_member_path"],
                        embedded_sidecars=sidecars,
                    )
                )
    except BadZipFile as exc:
        return _failed_container_result(metadata, f"ZIP container parse failed: {exc}")
    except ValueError as exc:
        return _failed_container_result(metadata, str(exc))

    if skipped_children:
        metadata["skipped_member_size_limit_count"] = len(skipped_children)
        metadata["warnings"] = sorted({*(metadata.get("warnings") or []), "member exceeds size limit"})
    return _final_container_result(metadata, children, direct_member_count=len(infos), total_uncompressed_bytes=total_bytes)


def _extract_tar_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    depth: int,
    member_prefix: str,
) -> ExtractionResult:
    format_name = "tar"
    metadata = _container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix)
    try:
        with tarfile.open(path) as archive:
            members = [member for member in archive.getmembers() if not member.isdir()]
            member_payloads: list[dict[str, Any]] = []
            skipped_children: list[ContainerChildAsset] = []
            total_bytes = 0
            for index, member in enumerate(members):
                member_path = _safe_container_member_name(member.name)
                if member_path is None:
                    return _failed_container_result(metadata, f"unsafe container member: {member.name}")
                if not member.isfile():
                    return _failed_container_result(metadata, f"unsafe non-file container member: {member_path}")
                size = int(member.size or 0)
                cap_message = _container_cap_message(policy, member_count=index + 1, member_size=0, total_bytes=total_bytes + size)
                if cap_message:
                    metadata["member_count"] = len(members)
                    return _metadata_only_container_result(metadata, cap_message)
                if size > policy.container_max_member_bytes:
                    skipped_children.append(
                        _skipped_container_child(
                            _join_container_member_path(member_prefix, member_path),
                            size_bytes=size,
                            reason="member_size_limit",
                            message="member exceeds size limit",
                            container_format=format_name,
                            container_kind=container_kind,
                            member_index=index,
                            compressed_size=None,
                            member_depth=depth + 1,
                            parent_member_path=member_prefix or None,
                        )
                    )
                    total_bytes += size
                    continue
                extracted = archive.extractfile(member)
                data = extracted.read() if extracted is not None else b""
                total_bytes += len(data)
                member_payloads.append(
                    {
                        "member_path": _join_container_member_path(member_prefix, member_path),
                        "data": data,
                        "member_index": index,
                        "compressed_size": None,
                        "member_depth": depth + 1,
                        "parent_member_path": member_prefix or None,
                    }
                )
            sidecars = _embedded_media_sidecars(member_payloads)
            children: list[ContainerChildAsset] = list(skipped_children)
            for member in member_payloads:
                children.extend(
                    _container_child_assets_from_bytes(
                        member["member_path"],
                        member["data"],
                        policy,
                        container_format=format_name,
                        container_kind=container_kind,
                        member_index=member["member_index"],
                        compressed_size=member["compressed_size"],
                        member_depth=member["member_depth"],
                        parent_member_path=member["parent_member_path"],
                        embedded_sidecars=sidecars,
                    )
                )
    except tarfile.TarError as exc:
        return _failed_container_result(metadata, f"TAR container parse failed: {exc}")
    except ValueError as exc:
        return _failed_container_result(metadata, str(exc))

    if skipped_children:
        metadata["skipped_member_size_limit_count"] = len(skipped_children)
        metadata["warnings"] = sorted({*(metadata.get("warnings") or []), "member exceeds size limit"})
    return _final_container_result(metadata, children, direct_member_count=len(members), total_uncompressed_bytes=total_bytes)


def _extract_stream_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    format_name: str,
    depth: int,
    member_prefix: str,
) -> ExtractionResult:
    metadata = _container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix)
    opener: Callable[..., Any]
    if format_name == "gzip":
        opener = gzip.open
    elif format_name == "bzip2":
        opener = bz2.open
    else:
        opener = lzma.open
    try:
        with opener(path, "rb") as handle:
            data = handle.read(policy.container_max_member_bytes + 1)
    except OSError as exc:
        return _failed_container_result(metadata, f"{format_name} stream parse failed: {exc}")
    if len(data) > policy.container_max_member_bytes:
        return _metadata_only_container_result(metadata, "member exceeds size limit")
    member_path = _stream_member_name(path, format_name)
    children = _container_child_assets_from_bytes(
        _join_container_member_path(member_prefix, member_path),
        data,
        policy,
        container_format=format_name,
        container_kind=container_kind,
        member_index=0,
        compressed_size=path.stat().st_size,
        member_depth=depth + 1,
        parent_member_path=member_prefix or None,
    )
    return _final_container_result(metadata, list(children), direct_member_count=1, total_uncompressed_bytes=len(data))


def _extract_tool_stream_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    format_name: str,
    depth: int,
    member_prefix: str,
) -> ExtractionResult | None:
    tool_name = "zstd" if format_name == "zst" else "lz4"
    command = shutil.which(tool_name)
    if command is None:
        return None
    metadata = _container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix, attempted=[tool_name])
    try:
        result = run_no_window(
            [command, "-dc", str(path)],
            capture_output=True,
            timeout=120,
            check=False,
        )
    except Exception as exc:
        return _failed_container_result(metadata, str(exc))
    if result.returncode != 0:
        return _failed_container_result(metadata, result.stderr.decode("utf-8", errors="replace") if isinstance(result.stderr, bytes) else str(result.stderr or "stream decompression failed"))
    data = bytes(result.stdout or b"")
    if len(data) > policy.container_max_member_bytes:
        return _metadata_only_container_result(metadata, "member exceeds size limit")
    children = _container_child_assets_from_bytes(
        _join_container_member_path(member_prefix, _stream_member_name(path, format_name)),
        data,
        policy,
        container_format=format_name,
        container_kind=container_kind,
        member_index=0,
        compressed_size=path.stat().st_size,
        member_depth=depth + 1,
        parent_member_path=member_prefix or None,
    )
    return _final_container_result(metadata, list(children), direct_member_count=1, total_uncompressed_bytes=len(data))


def _extract_optional_tool_container(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    format_name: str,
    depth: int,
    member_prefix: str,
) -> ExtractionResult | None:
    tools = OPTIONAL_CONTAINER_TOOLS.get(format_name, ())
    for tool in tools:
        command = shutil.which(tool)
        if command is None:
            continue
        result = _extract_with_directory_tool(
            path,
            policy,
            container_kind=container_kind,
            format_name=format_name,
            depth=depth,
            member_prefix=member_prefix,
            command=command,
            command_name=tool,
            attempted=list(tools),
        )
        if result is not None:
            return result
    return None


def _extract_with_directory_tool(
    path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    format_name: str,
    depth: int,
    member_prefix: str,
    command: str,
    command_name: str,
    attempted: list[str],
) -> ExtractionResult | None:
    metadata = _container_metadata(format_name, container_kind, policy=policy, depth=depth, member_prefix=member_prefix, attempted=attempted)
    with tempfile.TemporaryDirectory(prefix="flux-kb-container-") as temp_dir:
        temp_path = Path(temp_dir)
        command_line: list[str]
        kwargs: dict[str, Any] = {"text": True, "capture_output": True, "timeout": 180, "check": False}
        if command_name in {"7z", "7zz"}:
            command_line = [command, "x", "-y", f"-o{temp_dir}", str(path)]
        elif command_name == "bsdtar":
            command_line = [command, "-xf", str(path), "-C", temp_dir]
        elif command_name == "unar":
            command_line = [command, "-quiet", "-force-overwrite", "-output-directory", temp_dir, str(path)]
        elif command_name == "unrar":
            command_line = [command, "x", "-o+", "-inul", str(path), temp_dir]
        elif command_name == "ar":
            command_line = [command, "x", str(path)]
            kwargs["cwd"] = temp_dir
        else:
            return None
        try:
            result = run_no_window(command_line, **kwargs)
        except Exception as exc:
            return _failed_container_result(metadata, str(exc))
        if result.returncode != 0:
            return None
        return _children_from_extracted_directory(
            temp_path,
            policy,
            container_kind=container_kind,
            format_name=format_name,
            metadata=metadata,
            depth=depth,
            member_prefix=member_prefix,
        )
    return None


def _children_from_extracted_directory(
    temp_path: Path,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    format_name: str,
    metadata: dict[str, Any],
    depth: int,
    member_prefix: str,
) -> ExtractionResult:
    children: list[ContainerChildAsset] = []
    member_payloads: list[dict[str, Any]] = []
    total_bytes = 0
    files = [path for path in temp_path.rglob("*") if path.is_file() and not path.is_symlink()]
    for index, member in enumerate(files):
        member_path = _safe_container_member_name(member.relative_to(temp_path).as_posix())
        if member_path is None:
            return _failed_container_result(metadata, f"unsafe container member: {member}")
        data = member.read_bytes()
        cap_message = _container_cap_message(policy, member_count=index + 1, member_size=len(data), total_bytes=total_bytes + len(data))
        if cap_message:
            metadata["member_count"] = len(files)
            return _metadata_only_container_result(metadata, cap_message)
        total_bytes += len(data)
        member_payloads.append(
            {
                "member_path": _join_container_member_path(member_prefix, member_path),
                "data": data,
                "member_index": index,
                "compressed_size": None,
                "member_depth": depth + 1,
                "parent_member_path": member_prefix or None,
            }
        )
    sidecars = _embedded_media_sidecars(member_payloads)
    for member in member_payloads:
        children.extend(
            _container_child_assets_from_bytes(
                member["member_path"],
                member["data"],
                policy,
                container_format=format_name,
                container_kind=container_kind,
                member_index=member["member_index"],
                compressed_size=member["compressed_size"],
                member_depth=member["member_depth"],
                parent_member_path=member["parent_member_path"],
                embedded_sidecars=sidecars,
            )
        )
    return _final_container_result(metadata, children, direct_member_count=len(files), total_uncompressed_bytes=total_bytes)


def _container_child_from_bytes(
    member_path: str,
    data: bytes,
    policy: CorpusPolicy,
    *,
    container_format: str,
    container_kind: str,
    member_index: int,
    compressed_size: int | None,
    member_depth: int,
    parent_member_path: str | None,
) -> ContainerChildAsset:
    size = len(data)
    content_hash = hashlib.sha256(data).hexdigest()
    file_kind = _member_file_kind(member_path)
    extension = PurePosixPath(member_path).suffix.lower()
    mime_type, _ = mimetypes.guess_type(member_path)
    metadata: dict[str, Any] = {
        "extractor": "container_member",
        "container_format": container_format,
        "container_kind": container_kind,
        "container_member_path": member_path,
        "container_member_index": member_index,
        "container_depth": member_depth,
    }
    if parent_member_path:
        metadata["container_parent_path"] = parent_member_path
    if compressed_size is not None:
        metadata["compressed_size_bytes"] = compressed_size
    chunks: tuple[AssetChunk, ...] = ()
    extraction_tier = "metadata_only"
    extraction_status = "metadata_only"
    if file_kind in {"text", "code"} and size <= policy.max_inline_bytes:
        text = data.decode("utf-8", errors="replace").strip()
        chunks = _chunks_from_text(text, member_path)
        extraction_tier = "inline"
        extraction_status = "indexed" if chunks else "metadata_only"
    elif file_kind in {"archive", "container"}:
        metadata["nested_container"] = True
        if member_depth >= policy.container_max_depth:
            metadata["recursive_skipped_reason"] = "max_depth"
    return ContainerChildAsset(
        member_path=member_path,
        file_kind=file_kind,
        mime_type=mime_type,
        extension=extension,
        size_bytes=size,
        quick_hash=_container_child_quick_hash(member_path, size, content_hash),
        content_hash=content_hash,
        extraction_tier=extraction_tier,
        extraction_status=extraction_status,
        chunks=chunks,
        metadata=metadata,
    )


def _skipped_container_child(
    member_path: str,
    *,
    size_bytes: int,
    reason: str,
    message: str,
    container_format: str,
    container_kind: str,
    member_index: int,
    compressed_size: int | None,
    member_depth: int,
    parent_member_path: str | None,
) -> ContainerChildAsset:
    file_kind = _member_file_kind(member_path)
    extension = PurePosixPath(member_path).suffix.lower()
    mime_type, _ = mimetypes.guess_type(member_path)
    metadata: dict[str, Any] = {
        "extractor": "container_member",
        "container_format": container_format,
        "container_kind": container_kind,
        "container_member_path": member_path,
        "container_member_index": member_index,
        "container_depth": member_depth,
        "skipped_reason": reason,
        "warnings": [message],
    }
    if parent_member_path:
        metadata["container_parent_path"] = parent_member_path
    if compressed_size is not None:
        metadata["compressed_size_bytes"] = compressed_size
    return ContainerChildAsset(
        member_path=member_path,
        file_kind=file_kind,
        mime_type=mime_type,
        extension=extension,
        size_bytes=size_bytes,
        quick_hash=_container_child_quick_hash(member_path, size_bytes, reason),
        content_hash=None,
        extraction_tier="metadata_only",
        extraction_status="metadata_only",
        chunks=(),
        metadata=metadata,
    )


def _container_child_assets_from_bytes(
    member_path: str,
    data: bytes,
    policy: CorpusPolicy,
    *,
    container_format: str,
    container_kind: str,
    member_index: int,
    compressed_size: int | None,
    member_depth: int,
    parent_member_path: str | None,
    embedded_sidecars: dict[str, tuple[str, str]] | None = None,
) -> tuple[ContainerChildAsset, ...]:
    child = _container_child_from_bytes(
        member_path,
        data,
        policy,
        container_format=container_format,
        container_kind=container_kind,
        member_index=member_index,
        compressed_size=compressed_size,
        member_depth=member_depth,
        parent_member_path=parent_member_path,
    )
    children = [child]
    if child.file_kind in {"archive", "container"}:
        if member_depth >= policy.container_max_depth:
            return tuple(children)
        result = _extract_embedded_container_member(
            member_path,
            data,
            policy,
            container_kind=child.file_kind,
            depth=member_depth,
        )
        children[0] = _container_child_with_result(child, result, embedded=False)
        children.extend(result.child_assets)
    elif child.file_kind in EMBEDDED_PARSE_MEMBER_KINDS and member_depth <= policy.container_max_depth:
        result = _extract_embedded_member(
            member_path,
            data,
            policy,
            child.file_kind,
            sidecar=(embedded_sidecars or {}).get(member_path),
        )
        children[0] = _container_child_with_result(child, result, embedded=True)
    return tuple(children)


def _embedded_media_sidecars(member_payloads: list[dict[str, Any]]) -> dict[str, tuple[str, str]]:
    sidecars: dict[str, tuple[str, str]] = {}
    for member in member_payloads:
        sidecar_path = str(member["member_path"])
        lower_sidecar_path = sidecar_path.lower()
        for suffix in MEDIA_TRANSCRIPT_SIDECAR_SUFFIXES:
            if not lower_sidecar_path.endswith(suffix):
                continue
            media_member_path = sidecar_path[: -len(suffix)]
            if _member_file_kind(media_member_path) not in {"audio", "video"}:
                continue
            text = bytes(member["data"]).decode("utf-8", errors="replace").strip()
            if text:
                sidecars.setdefault(media_member_path, (sidecar_path, text))
    return sidecars


def _extract_embedded_container_member(
    member_path: str,
    data: bytes,
    policy: CorpusPolicy,
    *,
    container_kind: str,
    depth: int,
) -> ExtractionResult:
    with tempfile.TemporaryDirectory(prefix="flux-kb-member-") as temp_dir:
        temp_path = _materialize_member(Path(temp_dir), member_path, data)
        return _extract_container(
            temp_path,
            policy,
            container_kind=container_kind,
            depth=depth,
            member_prefix=member_path,
        )


def _extract_embedded_member(
    member_path: str,
    data: bytes,
    policy: CorpusPolicy,
    file_kind: str,
    *,
    sidecar: tuple[str, str] | None = None,
) -> ExtractionResult:
    try:
        with tempfile.TemporaryDirectory(prefix="flux-kb-member-") as temp_dir:
            temp_path = _materialize_member(Path(temp_dir), member_path, data)
            if file_kind == "document":
                return _extract_document(temp_path, policy)
            if file_kind == "diagram":
                return _extract_diagram(temp_path)
            if file_kind == "image":
                return _extract_image(temp_path)
            if file_kind in {"audio", "video"}:
                return _extract_media(temp_path, file_kind, embedded_sidecar=sidecar)
            if file_kind == "subtitle":
                return _extract_subtitle(temp_path)
            if file_kind == "mail":
                return _extract_mail(temp_path)
            if file_kind == "calendar":
                return _extract_calendar(temp_path)
            if file_kind == "contact":
                return _extract_contact(temp_path)
            if file_kind == "structured_data":
                return _extract_structured_data(temp_path, policy)
            if file_kind == "report":
                return _extract_report(temp_path)
            if file_kind == "database":
                return _extract_database(temp_path)
            if file_kind in {"geospatial", "cad", "scientific"}:
                return _extract_domain_metadata(temp_path, file_kind)
            if file_kind == "sensitive_metadata":
                return _extract_sensitive_metadata(temp_path)
    except Exception as exc:
        return ExtractionResult(
            status="failed",
            metadata={"extractor": file_kind, "warnings": [f"embedded member parse failed: {exc}"]},
            message=str(exc),
        )
    return ExtractionResult(status="metadata_only", metadata={"extractor": file_kind})


def _materialize_member(temp_dir: Path, member_path: str, data: bytes) -> Path:
    name = PurePosixPath(member_path).name or "member"
    target = temp_dir / name
    target.write_bytes(data)
    return target


def _container_child_with_result(
    child: ContainerChildAsset,
    result: ExtractionResult,
    *,
    embedded: bool,
) -> ContainerChildAsset:
    metadata = dict(child.metadata)
    metadata["embedded_extraction_status"] = result.status
    extractor = result.metadata.get("extractor") if isinstance(result.metadata, dict) else None
    if extractor:
        metadata["embedded_extractor"] = extractor
    warnings = result.metadata.get("warnings") if isinstance(result.metadata, dict) else None
    if warnings:
        metadata["warnings"] = list(warnings)
    if result.message:
        metadata["embedded_message"] = result.message
    if isinstance(result.metadata, dict):
        for key in ("transcript_source", "embedded_sidecar_path"):
            value = result.metadata.get(key)
            if value:
                metadata[key] = value
        for key in ("decorative", "vision", "frame_sampling"):
            summary = _embedded_visual_summary(result.metadata.get(key))
            if summary:
                metadata[f"embedded_{key}"] = summary
        for key in (
            "source_extension",
            "subtitle_format",
            "cue_count",
            "mail_format",
            "message_count",
            "attachment_count",
            "calendar_format",
            "event_count",
            "contact_format",
            "contact_count",
            "report_format",
            "finding_count",
            "test_count",
            "entry_count",
            "database_format",
            "table_count",
            "sensitive",
            "metadata_first",
            "tool_candidates",
        ):
            if key in result.metadata:
                metadata[f"embedded_{key}"] = result.metadata[key]
    if child.file_kind in {"archive", "container"}:
        metadata["nested_member_count"] = int(result.metadata.get("member_count") or 0)
        metadata["nested_parsed_child_count"] = int(result.metadata.get("parsed_child_count") or 0)
        metadata["nested_skipped_child_count"] = int(result.metadata.get("skipped_child_count") or 0)
        metadata["nested_blocked_dependency_count"] = int(result.metadata.get("blocked_dependency_count") or 0)
    status = result.status
    extraction_tier = "inline" if embedded and status == "indexed" else child.extraction_tier
    return replace(
        child,
        extraction_status=status,
        extraction_tier=extraction_tier,
        chunks=tuple(result.chunks or ()),
        metadata=metadata,
    )


def _embedded_visual_summary(value: Any) -> dict[str, Any] | None:
    if not isinstance(value, dict):
        return None
    allowed = {
        "status",
        "reason",
        "provider",
        "model",
        "cache_hits",
        "cache_misses",
        "descriptions",
        "blocked_dependency_count",
        "frame_count",
        "thumbnail_cache_hits",
        "thumbnail_cache_misses",
        "timestamps",
        "scene_scores",
    }
    summary = {key: value[key] for key in allowed if key in value}
    return summary or None


def _member_file_kind(member_path: str) -> str:
    path = PurePosixPath(member_path)
    name = path.name.lower()
    ext = path.suffix.lower()
    if ext in SENSITIVE_METADATA_EXTENSIONS:
        return "sensitive_metadata"
    if ext in SUBTITLE_EXTENSIONS:
        return "subtitle"
    if ext in MAIL_EXTENSIONS:
        return "mail"
    if ext in CALENDAR_EXTENSIONS:
        return "calendar"
    if ext in CONTACT_EXTENSIONS:
        return "contact"
    if ext in STRUCTURED_DATA_EXTENSIONS:
        return "structured_data"
    if ext in REPORT_EXTENSIONS or name in REPORT_NAMES:
        return "report"
    if ext in {".db", ".sqlite", ".sqlite3", ".duckdb", ".mdb", ".accdb", ".dbf", ".fdb"}:
        return "database"
    if ext in GEOSPATIAL_EXTENSIONS:
        return "geospatial"
    if ext in CAD_EXTENSIONS:
        return "cad"
    if ext in SCIENTIFIC_EXTENSIONS:
        return "scientific"
    if name in TEXT_MEMBER_NAMES:
        return "text"
    if ext in DIAGRAM_EXTENSIONS or any(name.endswith(suffix) for suffix in DIAGRAM_COMPOUND_SUFFIXES):
        return "diagram"
    if ext in CODE_EXTENSIONS:
        return "code"
    if ext in TEXT_MEMBER_EXTENSIONS:
        return "text"
    if ext in DOCUMENT_EXTENSIONS:
        return "document"
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in AUDIO_EXTENSIONS:
        return "audio"
    if ext in VIDEO_EXTENSIONS:
        return "video"
    if ext in CONTAINER_EXTENSIONS:
        return "container"
    if ext in ARCHIVE_EXTENSIONS or any(name.endswith(suffix) for suffix in ARCHIVE_COMPOUND_SUFFIXES):
        return "archive"
    mime_type, _ = mimetypes.guess_type(member_path)
    if mime_type and mime_type.startswith("text/"):
        return "text"
    return "binary"


def _container_metadata(
    format_name: str,
    container_kind: str,
    *,
    policy: CorpusPolicy,
    depth: int,
    member_prefix: str,
    attempted: list[str] | None = None,
) -> dict[str, Any]:
    metadata: dict[str, Any] = {
        "extractor": "container",
        "format": format_name,
        "container_kind": container_kind,
        "member_count": 0,
        "container_depth": depth,
        "max_depth": policy.container_max_depth,
        "parsed_child_count": 0,
        "skipped_child_count": 0,
        "blocked_dependency_count": 0,
    }
    if member_prefix:
        metadata["container_member_path"] = member_prefix
    if attempted is not None:
        metadata["attempted"] = attempted
    return metadata


def _final_container_result(
    metadata: dict[str, Any],
    children: list[ContainerChildAsset],
    *,
    direct_member_count: int,
    total_uncompressed_bytes: int,
) -> ExtractionResult:
    metadata["member_count"] = direct_member_count
    metadata["child_asset_count"] = len(children)
    metadata["total_uncompressed_bytes"] = total_uncompressed_bytes
    metadata["parsed_child_count"] = sum(1 for child in children if child.extraction_status == "indexed")
    metadata["blocked_dependency_count"] = sum(1 for child in children if child.extraction_status == "blocked_missing_dependency")
    metadata["skipped_child_count"] = sum(
        1 for child in children if child.extraction_status not in {"indexed", "blocked_missing_dependency"}
    )
    metadata.update(_container_visual_telemetry(children))
    return ExtractionResult(status="metadata_only", child_assets=tuple(children), metadata=metadata)


def _container_visual_telemetry(children: list[ContainerChildAsset]) -> dict[str, int]:
    telemetry = {
        "vision_cache_hits": 0,
        "vision_cache_misses": 0,
        "vision_descriptions": 0,
        "vision_blocked_dependency_count": 0,
        "decorative_image_skips": 0,
        "frame_sample_count": 0,
        "thumbnail_cache_hits": 0,
        "thumbnail_cache_misses": 0,
    }
    for child in children:
        decorative = child.metadata.get("embedded_decorative")
        if isinstance(decorative, dict) and decorative.get("status") == "skipped":
            telemetry["decorative_image_skips"] += 1
        vision = child.metadata.get("embedded_vision")
        if isinstance(vision, dict):
            telemetry["vision_cache_hits"] += int(vision.get("cache_hits") or 0)
            telemetry["vision_cache_misses"] += int(vision.get("cache_misses") or 0)
            telemetry["vision_descriptions"] += int(vision.get("descriptions") or 0)
            telemetry["vision_blocked_dependency_count"] += int(vision.get("blocked_dependency_count") or 0)
        frame_sampling = child.metadata.get("embedded_frame_sampling")
        if isinstance(frame_sampling, dict):
            telemetry["frame_sample_count"] += int(frame_sampling.get("frame_count") or 0)
            telemetry["thumbnail_cache_hits"] += int(frame_sampling.get("thumbnail_cache_hits") or 0)
            telemetry["thumbnail_cache_misses"] += int(frame_sampling.get("thumbnail_cache_misses") or 0)
    return telemetry


def _join_container_member_path(prefix: str, member_path: str) -> str:
    if not prefix:
        return member_path
    return f"{prefix.rstrip('/')}/{member_path.lstrip('/')}"


def _container_cap_message(policy: CorpusPolicy, *, member_count: int, member_size: int, total_bytes: int) -> str | None:
    if policy.container_max_depth < 1:
        return "container expansion depth must be at least 1"
    if member_count > policy.container_max_members:
        return "container exceeds member count limit"
    if member_size > policy.container_max_member_bytes:
        return "member exceeds size limit"
    if total_bytes > policy.container_max_total_bytes:
        return "container exceeds total size limit"
    return None


def _metadata_only_container_result(metadata: dict[str, Any], message: str) -> ExtractionResult:
    metadata = {**metadata, "warnings": [message]}
    return ExtractionResult(status="metadata_only", metadata=metadata, message=message)


def _failed_container_result(metadata: dict[str, Any], message: str) -> ExtractionResult:
    metadata = {**metadata, "warnings": [message]}
    return ExtractionResult(status="failed", metadata=metadata, message=message)


def _safe_container_member_name(name: str) -> str | None:
    normalized = str(name or "").replace("\\", "/")
    if not normalized or normalized.startswith("/"):
        return None
    path = PurePosixPath(normalized)
    if path.is_absolute():
        return None
    if any(part in {"", ".", ".."} or ":" in part for part in path.parts):
        return None
    return path.as_posix()


def _stream_member_name(path: Path, format_name: str) -> str:
    name = path.name
    suffix_by_format = {"gzip": ".gz", "bzip2": ".bz2", "xz": ".xz", "zst": ".zst", "lz4": ".lz4"}
    suffix = suffix_by_format.get(format_name, path.suffix)
    if name.lower().endswith(suffix):
        candidate = name[: -len(suffix)]
        if candidate:
            return candidate
    return path.stem or f"{path.name}.contents"


def _container_child_quick_hash(member_path: str, size_bytes: int, content_hash: str) -> str:
    value = f"{member_path}:{size_bytes}:{content_hash}".encode("utf-8", errors="ignore")
    return hashlib.sha256(value).hexdigest()


def _extract_diagram(path: Path) -> ExtractionResult:
    name = path.name.lower()
    if path.suffix.lower() in {".vsdm", ".vsdx", ".vssm", ".vssx", ".vstm", ".vstx"}:
        return _extract_vsdx_diagram(path)
    if path.suffix.lower() in {".dio", ".drawio"} or name.endswith((".drawio.svg", ".drawio.png")):
        return _extract_drawio_diagram(path)
    return ExtractionResult(status="metadata_only", metadata={"extractor": "diagram", "format": "unknown"})


def _extract_drawio_diagram(path: Path) -> ExtractionResult:
    metadata = _diagram_metadata("drawio")
    try:
        raw = _read_limited_file(path, DIAGRAM_MAX_TOTAL_BYTES)
        mxfile_xml = _embedded_mxfile_xml(raw)
        if not mxfile_xml:
            metadata["warnings"].append("drawio mxfile payload not found")
            return ExtractionResult(status="metadata_only", metadata=metadata)
        root = ElementTree.fromstring(mxfile_xml)
    except ValueError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc))
    except ElementTree.ParseError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=f"drawio XML parse failed: {exc}")

    page_blocks: list[str] = []
    try:
        pages = _drawio_pages(root)
    except ValueError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc))
    metadata["page_count"] = len(pages)
    for page_name, page_root in pages:
        if page_root is None:
            metadata["warnings"].append(f"page payload could not be decoded: {page_name}")
            continue
        page_lines, page_stats = _summarize_drawio_page(page_name, page_root)
        _merge_diagram_stats(metadata, page_stats)
        if page_lines:
            page_blocks.append("\n".join([f"Page: {page_name}", *page_lines]))

    text = "\n\n".join(page_blocks)
    chunks = _chunks_from_text(text, path.name, modality="diagram")
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata=metadata,
    )


def _extract_vsdx_diagram(path: Path) -> ExtractionResult:
    metadata = _diagram_metadata("vsdx")
    try:
        with ZipFile(path) as archive:
            infos = archive.infolist()
            if len(infos) > DIAGRAM_MAX_ZIP_MEMBERS:
                raise ValueError(f"diagram container has too many members: {len(infos)}")
            total_page_xml_bytes = 0
            page_members = []
            for info in infos:
                member_name = _safe_zip_member_name(info.filename)
                if member_name is None:
                    raise ValueError(f"unsafe diagram container member: {info.filename}")
                normalized = member_name.lower()
                if normalized.startswith("visio/pages/") and normalized.endswith(".xml"):
                    if info.file_size > DIAGRAM_MAX_MEMBER_BYTES:
                        raise ValueError(f"diagram page XML member exceeds size limit: {member_name}")
                    total_page_xml_bytes += int(info.file_size or 0)
                    if total_page_xml_bytes > DIAGRAM_MAX_TOTAL_BYTES:
                        raise ValueError("diagram page XML exceeds readable XML limit")
                    page_members.append((member_name, info))
                    if len(page_members) > DIAGRAM_MAX_PAGE_XML_MEMBERS:
                        raise ValueError(f"diagram container has too many page XML members: {len(page_members)}")

            page_blocks: list[str] = []
            metadata["page_count"] = len(page_members)
            for member_name, info in page_members:
                xml_text = archive.read(info).decode("utf-8", errors="replace")
                page_name = PurePosixPath(member_name).stem
                try:
                    page_root = ElementTree.fromstring(xml_text)
                except ElementTree.ParseError as exc:
                    metadata["warnings"].append(f"{member_name}: XML parse failed: {exc}")
                    continue
                page_lines, page_stats = _summarize_vsdx_page(page_name, page_root)
                _merge_diagram_stats(metadata, page_stats)
                if page_lines:
                    page_blocks.append("\n".join([f"Page: {page_name}", *page_lines]))
    except BadZipFile as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=f"vsdx ZIP parse failed: {exc}")
    except ValueError as exc:
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc))

    if metadata["page_count"] == 0:
        metadata["warnings"].append("no Visio page XML members found")
    text = "\n\n".join(page_blocks)
    chunks = _chunks_from_text(text, path.name, modality="diagram")
    return ExtractionResult(
        status="indexed" if chunks else "metadata_only",
        chunks=chunks,
        metadata=metadata,
    )


def _diagram_metadata(format_name: str) -> dict[str, Any]:
    return {
        "extractor": "diagram",
        "format": format_name,
        "page_count": 0,
        "shape_count": 0,
        "connector_count": 0,
        "text_count": 0,
        "warnings": [],
    }


def _read_limited_file(path: Path, limit: int) -> bytes:
    size = path.stat().st_size
    if size > limit:
        raise ValueError("diagram file exceeds readable size limit")
    return path.read_bytes()


def _embedded_mxfile_xml(raw: bytes) -> str:
    text = raw.decode("utf-8", errors="ignore")
    for candidate in (text, unescape(text)):
        start = candidate.find("<mxfile")
        end = candidate.find("</mxfile>")
        if start >= 0 and end >= start:
            return candidate[start : end + len("</mxfile>")]
        if candidate.lstrip().startswith("<mxfile"):
            return candidate.strip()
    return ""


def _drawio_pages(root: ElementTree.Element) -> list[tuple[str, ElementTree.Element | None]]:
    root_name = _local_name(root.tag)
    if root_name == "mxGraphModel":
        return [("Page 1", root)]
    diagrams = [element for element in root.iter() if _local_name(element.tag) == "diagram"]
    pages: list[tuple[str, ElementTree.Element | None]] = []
    for index, diagram in enumerate(diagrams, start=1):
        page_name = _clean_diagram_text(diagram.attrib.get("name") or f"Page {index}") or f"Page {index}"
        child_models = [child for child in list(diagram) if _local_name(child.tag) == "mxGraphModel"]
        if child_models:
            pages.append((page_name, child_models[0]))
            continue
        decoded = _decode_drawio_payload("".join(diagram.itertext()).strip())
        if decoded:
            try:
                pages.append((page_name, ElementTree.fromstring(decoded)))
            except ElementTree.ParseError:
                pages.append((page_name, None))
        else:
            pages.append((page_name, None))
    return pages


def _decode_drawio_payload(value: str) -> str:
    text = value.strip()
    if not text:
        return ""
    candidates = [unescape(text), unquote(unescape(text))]
    for encoded in {text, unquote(unescape(text))}:
        try:
            compressed = base64.b64decode("".join(encoded.split()), validate=True)
        except (binascii.Error, ValueError):
            continue
        if len(compressed) > DIAGRAM_MAX_MEMBER_BYTES:
            raise ValueError("embedded drawio payload exceeds size limit")
        for wbits in (-15, zlib.MAX_WBITS):
            try:
                inflated = zlib.decompress(compressed, wbits)
            except zlib.error:
                continue
            if len(inflated) > DIAGRAM_MAX_MEMBER_BYTES:
                raise ValueError("inflated drawio payload exceeds size limit")
            decoded = inflated.decode("utf-8", errors="replace")
            candidates.extend([decoded, unquote(decoded)])
    for candidate in candidates:
        cleaned = candidate.strip()
        if cleaned.startswith("<"):
            return cleaned
    return ""


def _summarize_drawio_page(page_name: str, root: ElementTree.Element) -> tuple[list[str], dict[str, int]]:
    lines: list[str] = []
    stats = {"shape_count": 0, "connector_count": 0, "text_count": 0}
    for cell in root.iter():
        if _local_name(cell.tag) != "mxCell":
            continue
        value = _clean_diagram_text(cell.attrib.get("value"))
        link = _clean_diagram_text(cell.attrib.get("link"))
        source = _clean_diagram_text(cell.attrib.get("source"))
        target = _clean_diagram_text(cell.attrib.get("target"))
        is_connector = cell.attrib.get("edge") == "1" or bool(source or target)
        if is_connector:
            stats["connector_count"] += 1
            parts = []
            if value:
                parts.append(value)
                stats["text_count"] += 1
            if source or target:
                parts.append(f"{source or '?'} -> {target or '?'}")
            if link:
                parts.append(f"link {link}")
                stats["text_count"] += 1
            if parts:
                lines.append(f"Connector: {'; '.join(parts)}")
            continue
        if value or link:
            stats["shape_count"] += 1
            parts = []
            if value:
                parts.append(value)
                stats["text_count"] += 1
            if link:
                parts.append(f"link {link}")
                stats["text_count"] += 1
            lines.append(f"Shape: {'; '.join(parts)}")
    return lines, stats


def _summarize_vsdx_page(page_name: str, root: ElementTree.Element) -> tuple[list[str], dict[str, int]]:
    lines: list[str] = []
    stats = {"shape_count": 0, "connector_count": 0, "text_count": 0}
    for shape in root.iter():
        if _local_name(shape.tag) != "Shape":
            continue
        label = _shape_label(shape)
        text = _shape_text(shape)
        if not text:
            continue
        stats["shape_count"] += 1
        stats["text_count"] += 1
        lines.append(f"Shape {label}: {text}" if label else f"Shape: {text}")
    for connect in root.iter():
        if _local_name(connect.tag) != "Connect":
            continue
        stats["connector_count"] += 1
        source = _clean_diagram_text(connect.attrib.get("FromSheet"))
        target = _clean_diagram_text(connect.attrib.get("ToSheet"))
        if source or target:
            lines.append(f"Connector: {source or '?'} -> {target or '?'}")
    return lines, stats


def _shape_label(shape: ElementTree.Element) -> str:
    return _clean_diagram_text(
        shape.attrib.get("NameU")
        or shape.attrib.get("Name")
        or shape.attrib.get("ID")
        or shape.attrib.get("id")
    )


def _shape_text(shape: ElementTree.Element) -> str:
    parts: list[str] = []
    for element in shape.iter():
        if _local_name(element.tag) == "Text":
            text = _clean_diagram_text(" ".join(element.itertext()))
            if text:
                parts.append(text)
    return " ".join(dict.fromkeys(parts))


def _merge_diagram_stats(metadata: dict[str, Any], stats: dict[str, int]) -> None:
    for key in ("shape_count", "connector_count", "text_count"):
        metadata[key] = int(metadata.get(key) or 0) + int(stats.get(key) or 0)


def _safe_zip_member_name(name: str) -> str | None:
    normalized = str(name or "").replace("\\", "/")
    if not normalized or normalized.startswith("/"):
        return None
    path = PurePosixPath(normalized)
    if path.is_absolute() or any(part in {"", ".", ".."} for part in path.parts):
        return None
    return path.as_posix()


def _clean_diagram_text(value: Any) -> str:
    text = unquote(unescape(str(value or "")))
    text = re.sub(r"<[^>]+>", " ", text)
    return " ".join(text.split())


def _local_name(tag: str) -> str:
    return str(tag).rsplit("}", 1)[-1]


def _extract_image(path: Path) -> ExtractionResult:
    metadata = image_metadata(path)
    decorative = _decorative_image_metadata(metadata)
    if decorative is not None:
        metadata["decorative"] = decorative
        return ExtractionResult(status="indexed", metadata=metadata)
    ocr = _ocr_image(path)
    metadata["ocr"] = ocr.metadata
    vision = _vision_image(path, source_label=path.name)
    if vision.status != "disabled":
        metadata["vision"] = vision.metadata
        metadata["vision_escalation"] = _vision_escalation_status(vision)
    chunks = _reindexed_chunks(
        _chunks_from_text(ocr.text, path.name, modality="ocr"),
        _chunks_from_text(vision.text, path.name, modality="vision"),
    )
    if chunks:
        return ExtractionResult(status="indexed", chunks=chunks, metadata=metadata)
    if ocr.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=ocr.message)
    if ocr.status == "blocked_missing_dependency":
        if vision.status in {"blocked_config", "blocked_missing_dependency"}:
            return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=vision.message)
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=ocr.message)
    if vision.status in {"blocked_config", "blocked_missing_dependency"}:
        if ocr.status == "completed":
            metadata["vision_escalation"] = "unavailable"
            return ExtractionResult(status="metadata_only", metadata=metadata, message=vision.message)
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=vision.message)
    if vision.status == "failed":
        if ocr.status == "completed":
            metadata["vision_escalation"] = "unavailable"
            return ExtractionResult(status="metadata_only", metadata=metadata, message=vision.message)
        return ExtractionResult(status="failed", metadata=metadata, message=vision.message)
    return ExtractionResult(status="metadata_only", metadata=metadata)


def _vision_escalation_status(vision: VisionResult) -> str:
    status = str(vision.metadata.get("status") or vision.status)
    if vision.status == "completed":
        return "completed" if vision.text.strip() else "no_content"
    if status.startswith("skipped_"):
        return "ineligible"
    if vision.status in {"blocked_config", "blocked_missing_dependency", "failed"}:
        return "unavailable"
    return status

def plan_staged_media_extraction(path: str | Path, file_kind: str | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    kind = str(file_kind or ("video" if file_path.suffix.lower() in VIDEO_EXTENSIONS else "audio"))
    metadata: dict[str, Any] = {"extractor": kind}
    sidecar = _read_sidecar_transcript(file_path)
    if sidecar:
        chunks = _chunks_from_text(sidecar, file_path.name, modality="transcript")
        return ExtractionResult(status="indexed", chunks=chunks, metadata={**metadata, "transcript_source": "sidecar"})
    probe_error, probed = _probe_media(file_path, metadata)
    if probe_error is not None:
        return probe_error
    duration = _media_duration_seconds(probed)
    first_job: dict[str, Any] | None = None
    followup_jobs: list[dict[str, Any]] = []
    frame_job = _planned_video_frame_job(file_path, probed) if kind == "video" else None
    if frame_job is not None:
        metadata["frame_sampling"] = frame_job["metadata"]
        followup_jobs.append({"job_type": frame_job["job_type"], "payload": frame_job["payload"]})
    elif kind == "video":
        frame_settings = _frame_sampling_settings()
        frame_status = "disabled" if not frame_settings["enabled"] else "skipped_no_duration"
        metadata["frame_sampling"] = _frame_sampling_metadata(
            status=frame_status,
            duration_seconds=duration,
            max_duration_seconds=frame_settings["max_duration_seconds"],
            frame_sample_count=frame_settings["frame_sample_count"],
            scene_threshold=frame_settings["scene_threshold"],
        )
    asr_job, asr_metadata, asr_message = _planned_media_asr_job(probed, kind)
    metadata["asr"] = asr_metadata
    if asr_job is not None:
        if followup_jobs:
            asr_job["payload"]["followup_jobs"] = followup_jobs
        first_job = asr_job
    elif followup_jobs:
        first_job = followup_jobs[0]
    if first_job is None:
        return ExtractionResult(status="metadata_only", metadata=metadata, message=asr_message)
    pending_job_count = _estimated_staged_media_job_count(first_job, duration)
    metadata["staged_jobs"] = [first_job]
    metadata["staged_extraction"] = {
        "status": "planned",
        "content_extracted": False,
        "pending_job_count": pending_job_count,
        "next_job_type": first_job["job_type"],
        "unit": "media",
        "duration_seconds": duration,
    }
    return ExtractionResult(status="staged", metadata=metadata)


def extract_media_segment(path: str | Path, payload: dict[str, Any] | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    payload = payload or {}
    kind = str(payload.get("file_kind") or ("video" if file_path.suffix.lower() in VIDEO_EXTENSIONS else "audio"))
    metadata: dict[str, Any] = {"extractor": "media_segment", "file_kind": kind}
    probe_error, probed = _probe_media(file_path, metadata)
    if probe_error is not None:
        return probe_error
    total_duration = _float_or_none(payload.get("duration_seconds"))
    if total_duration is None:
        total_duration = _media_duration_seconds(probed)
    segment_index = max(0, _optional_int(payload.get("segment_index")) or 0)
    segment_start = max(0.0, _float_or_none(payload.get("segment_start_seconds")) or 0.0)
    segment_duration = _float_or_none(payload.get("segment_duration_seconds"))
    asr = _asr_media(
        file_path,
        probed,
        segment_start_seconds=segment_start,
        segment_duration_seconds=segment_duration,
        segment_index=segment_index,
    )
    metadata["asr"] = asr.metadata
    if asr.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=asr.message)
    if asr.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=asr.message)
    local_chunks = _chunks_from_text(asr.text, f"{file_path.name} segment {segment_index + 1}", modality="transcript")
    chunks = _offset_chunks(
        local_chunks,
        segment_index * MEDIA_SEGMENT_CHUNK_INDEX_STRIDE,
        locator_prefix=f"media:{segment_start:g}-{_segment_end_label(segment_start, segment_duration)}",
        metadata={"media_segment_index": segment_index, "segment_start_seconds": segment_start},
    )
    chunks_seen = max(0, _optional_int(payload.get("chunks_seen")) or 0) + len(chunks)
    followup_jobs = _normalised_followup_jobs(payload.get("followup_jobs"))
    next_job = _next_media_segment_job(
        payload,
        total_duration=total_duration,
        segment_index=segment_index,
        segment_start=segment_start,
        segment_duration=segment_duration,
        chunks_seen=chunks_seen,
        followup_jobs=followup_jobs,
    )
    metadata["staged_extraction"] = {
        "status": "piece_completed",
        "complete": next_job is None,
        "unit": "media_segment",
        "segment_index": segment_index,
        "segment_start_seconds": segment_start,
        "segment_duration_seconds": segment_duration,
        "duration_seconds": total_duration,
        "chunks_written": len(chunks),
        "chunks_seen": chunks_seen,
    }
    if next_job is not None:
        metadata["staged_extraction"]["next_job"] = next_job
        return ExtractionResult(status="staged", chunks=chunks, metadata=metadata, message=asr.message)
    return ExtractionResult(
        status="indexed" if chunks_seen > 0 else "metadata_only",
        chunks=chunks,
        metadata=metadata,
        message=asr.message,
    )


def extract_video_frames(path: str | Path, payload: dict[str, Any] | None = None) -> ExtractionResult:
    file_path = _extractor_path(path)
    payload = payload or {}
    metadata: dict[str, Any] = {"extractor": "video_frames", "file_kind": "video"}
    probe_error, probed = _probe_media(file_path, metadata)
    if probe_error is not None:
        return probe_error
    raw_timestamps = payload.get("timestamps")
    timestamps = [float(value) for value in raw_timestamps] if isinstance(raw_timestamps, list) else []
    if not timestamps:
        timestamps = _planned_frame_timestamps(_media_duration_seconds(probed), _frame_sampling_settings()["frame_sample_count"])
    frame_sampling = _sample_video_frames_at_timestamps(file_path, probed, timestamps=timestamps)
    metadata["frame_sampling"] = frame_sampling.metadata
    vision_summary = _vision_summary_from_frame_sampling(frame_sampling)
    if vision_summary is not None:
        metadata["vision"] = vision_summary
    if frame_sampling.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=frame_sampling.message)
    if frame_sampling.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=frame_sampling.message)
    chunks = _offset_chunks(
        frame_sampling.chunks,
        VIDEO_FRAME_CHUNK_INDEX_BASE,
        metadata={"video_frame_sampling": True},
    )
    chunks_seen = max(0, _optional_int(payload.get("chunks_seen")) or 0) + len(chunks)
    followup_jobs = _normalised_followup_jobs(payload.get("followup_jobs"))
    next_job = followup_jobs[0] if followup_jobs else None
    if next_job is not None:
        remaining = followup_jobs[1:]
        payload_value = dict(next_job.get("payload") or {})
        payload_value["chunks_seen"] = chunks_seen
        if remaining:
            payload_value["followup_jobs"] = remaining
        next_job = {"job_type": str(next_job.get("job_type") or ""), "payload": payload_value}
    metadata["staged_extraction"] = {
        "status": "piece_completed",
        "complete": next_job is None,
        "unit": "video_frames",
        "chunks_written": len(chunks),
        "chunks_seen": chunks_seen,
    }
    if next_job is not None:
        metadata["staged_extraction"]["next_job"] = next_job
        return ExtractionResult(status="staged", chunks=chunks, metadata=metadata, message=frame_sampling.message)
    return ExtractionResult(
        status="indexed" if chunks_seen > 0 else "metadata_only",
        chunks=chunks,
        metadata=metadata,
        message=frame_sampling.message,
    )


def _probe_media(path: Path, metadata: dict[str, Any]) -> tuple[ExtractionResult | None, dict[str, Any]]:
    ffprobe = shutil.which("ffprobe")
    if ffprobe is None:
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message="ffprobe command not found"), {}
    try:
        result = run_no_window(
            [
                ffprobe,
                "-v",
                "error",
                "-show_format",
                "-show_streams",
                "-of",
                "json",
                str(path),
            ],
            text=True,
            capture_output=True,
            timeout=30,
            check=False,
        )
    except Exception as exc:  # pragma: no cover - environment-specific
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc)), {}
    if result.returncode != 0:
        return ExtractionResult(status="failed", metadata=metadata, message=result.stderr.strip() or "ffprobe failed"), {}
    try:
        probed = json.loads(result.stdout or "{}")
    except json.JSONDecodeError:
        probed = {}
    metadata["ffprobe"] = probed
    return None, probed


def _planned_media_asr_job(probed: dict[str, Any], file_kind: str) -> tuple[dict[str, Any] | None, dict[str, Any], str | None]:
    settings = _asr_settings()
    duration = _media_duration_seconds(probed)
    provider = str(settings["provider"] or "local_faster_whisper")
    metadata = _asr_metadata(
        status="planned",
        duration_seconds=duration,
        max_duration_seconds=settings["max_duration_seconds"],
        device=settings["device"],
        compute_type=settings["compute_type"],
        provider=provider,
        model=str(settings.get("model") or ""),
        base_url=str(settings.get("base_url") or ""),
    )
    has_audio_stream = _media_has_audio_stream(probed)
    if has_audio_stream is not None:
        metadata["has_audio_stream"] = has_audio_stream
    if not settings["enabled"]:
        return None, {**metadata, "status": "disabled"}, None
    if has_audio_stream is False:
        return None, {**metadata, "status": "skipped_no_audio_stream"}, "media has no audio stream for ASR"
    segment_duration = _media_asr_segment_seconds(settings, duration)
    payload: dict[str, Any] = {
        "file_kind": file_kind,
        "segment_index": 0,
        "segment_start_seconds": 0.0,
        "duration_seconds": duration,
        "chunks_seen": 0,
    }
    if segment_duration is not None:
        payload["segment_duration_seconds"] = segment_duration
    return {"job_type": "corpus_extract_media_segment", "payload": payload}, metadata, None


def _planned_video_frame_job(_path: Path, probed: dict[str, Any]) -> dict[str, Any] | None:
    settings = _frame_sampling_settings()
    duration = _media_duration_seconds(probed)
    metadata = _frame_sampling_metadata(
        status="planned",
        duration_seconds=duration,
        max_duration_seconds=settings["max_duration_seconds"],
        frame_sample_count=settings["frame_sample_count"],
        scene_threshold=settings["scene_threshold"],
    )
    if not settings["enabled"]:
        return None
    timestamps = _planned_frame_timestamps(duration, settings["frame_sample_count"])
    if not timestamps:
        return None
    metadata.update(
        {
            "status": "planned",
            "sampling_strategy": "staged_evenly_spaced",
            "timestamps": timestamps,
            "frame_count": len(timestamps),
        }
    )
    return {
        "job_type": "corpus_extract_video_frames",
        "payload": {
            "file_kind": "video",
            "timestamps": timestamps,
            "duration_seconds": duration,
            "chunks_seen": 0,
        },
        "metadata": metadata,
    }


def _media_asr_segment_seconds(settings: dict[str, Any], duration: float | None) -> float | None:
    if duration is None:
        return None
    configured = _optional_int(settings.get("max_duration_seconds")) or MEDIA_ASR_SEGMENT_SECONDS
    segment_seconds = max(60, min(MEDIA_ASR_SEGMENT_SECONDS, configured))
    return min(float(segment_seconds), max(0.0, duration))


def _planned_frame_timestamps(duration: float | None, count: int) -> list[float]:
    if duration is None or duration <= 0 or count <= 0:
        return []
    if count == 1:
        return [round(max(0.0, duration / 2), 3)]
    step = duration / (count + 1)
    return [round(max(0.0, min(duration, step * index)), 3) for index in range(1, count + 1)]


def _estimated_staged_media_job_count(first_job: dict[str, Any], duration: float | None) -> int:
    count = 1
    payload = first_job.get("payload") if isinstance(first_job.get("payload"), dict) else {}
    segment_duration = _float_or_none(payload.get("segment_duration_seconds")) if isinstance(payload, dict) else None
    if first_job.get("job_type") == "corpus_extract_media_segment" and duration is not None and segment_duration and segment_duration > 0:
        count = max(1, math.ceil(duration / segment_duration))
    followups = _normalised_followup_jobs(payload.get("followup_jobs") if isinstance(payload, dict) else None)
    return count + len(followups)


def _normalised_followup_jobs(value: Any) -> list[dict[str, Any]]:
    if not isinstance(value, list):
        return []
    jobs: list[dict[str, Any]] = []
    for item in value:
        if not isinstance(item, dict):
            continue
        job_type = str(item.get("job_type") or "").strip()
        if not job_type:
            continue
        payload = item.get("payload") if isinstance(item.get("payload"), dict) else {}
        jobs.append({"job_type": job_type, "payload": dict(payload)})
    return jobs


def _next_media_segment_job(
    payload: dict[str, Any],
    *,
    total_duration: float | None,
    segment_index: int,
    segment_start: float,
    segment_duration: float | None,
    chunks_seen: int,
    followup_jobs: list[dict[str, Any]],
) -> dict[str, Any] | None:
    if total_duration is not None and segment_duration is not None and segment_duration > 0:
        next_start = segment_start + segment_duration
        if next_start < total_duration - 0.001:
            next_payload = dict(payload)
            next_payload.update(
                {
                    "segment_index": segment_index + 1,
                    "segment_start_seconds": round(next_start, 3),
                    "segment_duration_seconds": min(segment_duration, total_duration - next_start),
                    "duration_seconds": total_duration,
                    "chunks_seen": chunks_seen,
                }
            )
            if followup_jobs:
                next_payload["followup_jobs"] = followup_jobs
            return {"job_type": "corpus_extract_media_segment", "payload": next_payload}
    if followup_jobs:
        next_job = followup_jobs[0]
        remaining = followup_jobs[1:]
        next_payload = dict(next_job.get("payload") or {})
        next_payload["chunks_seen"] = chunks_seen
        if remaining:
            next_payload["followup_jobs"] = remaining
        return {"job_type": str(next_job.get("job_type") or ""), "payload": next_payload}
    return None


def _segment_end_label(segment_start: float, segment_duration: float | None) -> str:
    if segment_duration is None:
        return "end"
    return f"{segment_start + segment_duration:g}"


def _extract_media(path: Path, file_kind: str, *, embedded_sidecar: tuple[str, str] | None = None) -> ExtractionResult:
    metadata: dict[str, Any] = {"extractor": file_kind}
    if embedded_sidecar is not None:
        sidecar_path, sidecar_text = embedded_sidecar
        chunks = _chunks_from_text(sidecar_text, path.name, modality="transcript")
        return ExtractionResult(
            status="indexed" if chunks else "metadata_only",
            chunks=chunks,
            metadata={
                **metadata,
                "transcript_source": "embedded_sidecar",
                "embedded_sidecar_path": sidecar_path,
            },
        )
    sidecar = _read_sidecar_transcript(path)
    if sidecar:
        chunks = _chunks_from_text(sidecar, path.name, modality="transcript")
        return ExtractionResult(status="indexed", chunks=chunks, metadata={**metadata, "transcript_source": "sidecar"})
    ffprobe = shutil.which("ffprobe")
    if ffprobe is None:
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message="ffprobe command not found")
    try:
        result = run_no_window(
            [
                ffprobe,
                "-v",
                "error",
                "-show_format",
                "-show_streams",
                "-of",
                "json",
                str(path),
            ],
            text=True,
            capture_output=True,
            timeout=30,
            check=False,
        )
    except Exception as exc:  # pragma: no cover - environment-specific
        return ExtractionResult(status="failed", metadata=metadata, message=str(exc))
    if result.returncode != 0:
        return ExtractionResult(status="failed", metadata=metadata, message=result.stderr.strip() or "ffprobe failed")
    try:
        probed = json.loads(result.stdout or "{}")
    except json.JSONDecodeError:
        probed = {}
    metadata["ffprobe"] = probed
    frame_sampling = _sample_video_frames(path, probed) if file_kind == "video" else None
    frame_chunks: tuple[AssetChunk, ...] = ()
    if frame_sampling is not None and frame_sampling.status != "disabled":
        metadata["frame_sampling"] = frame_sampling.metadata
        frame_chunks = frame_sampling.chunks
        vision_summary = _vision_summary_from_frame_sampling(frame_sampling)
        if vision_summary is not None:
            metadata["vision"] = vision_summary
    asr = _asr_media(path, probed)
    metadata["asr"] = asr.metadata
    chunks = _reindexed_chunks(_chunks_from_text(asr.text, path.name, modality="transcript"), frame_chunks)
    if chunks:
        return ExtractionResult(status="indexed", chunks=chunks, metadata=metadata, message=asr.message or frame_sampling.message if frame_sampling else asr.message)
    if asr.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=asr.message)
    if asr.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=asr.message)
    if frame_sampling is not None and frame_sampling.status == "blocked_missing_dependency":
        return ExtractionResult(status="blocked_missing_dependency", metadata=metadata, message=frame_sampling.message)
    if frame_sampling is not None and frame_sampling.status == "failed":
        return ExtractionResult(status="failed", metadata=metadata, message=frame_sampling.message)
    return ExtractionResult(
        status="metadata_only",
        metadata=metadata,
        message=asr.message or frame_sampling.message if frame_sampling else asr.message,
    )


def _asr_media(
    path: Path,
    probed: dict[str, Any],
    *,
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
    segment_index: int | None = None,
) -> AsrResult:
    settings = _asr_settings()
    duration = _media_duration_seconds(probed)
    provider = str(settings["provider"] or "local_faster_whisper")
    metadata = _asr_metadata(
        status="pending",
        duration_seconds=duration,
        max_duration_seconds=settings["max_duration_seconds"],
        device=settings["device"],
        compute_type=settings["compute_type"],
        provider=provider,
        model=str(settings.get("model") or ""),
        base_url=str(settings.get("base_url") or ""),
    )
    has_audio_stream = _media_has_audio_stream(probed)
    if has_audio_stream is not None:
        metadata["has_audio_stream"] = has_audio_stream
    if segment_index is not None:
        metadata["segment_index"] = int(segment_index)
    if segment_start_seconds is not None:
        metadata["segment_start_seconds"] = float(segment_start_seconds)
    if segment_duration_seconds is not None:
        metadata["segment_duration_seconds"] = float(segment_duration_seconds)
    if not settings["enabled"]:
        return AsrResult(status="completed", metadata={**metadata, "status": "disabled"})
    if has_audio_stream is False:
        return AsrResult(
            status="completed",
            metadata={**metadata, "status": "skipped_no_audio_stream"},
            message="media has no audio stream for ASR",
        )
    if provider == "openai_compatible":
        return _asr_media_openai_compatible(
            path,
            settings=settings,
            metadata=metadata,
            segment_start_seconds=segment_start_seconds,
            segment_duration_seconds=segment_duration_seconds,
        )
    return _asr_media_local_faster_whisper(
        path,
        settings=settings,
        metadata=metadata,
        segment_start_seconds=segment_start_seconds,
        segment_duration_seconds=segment_duration_seconds,
    )


def _asr_media_local_faster_whisper(
    path: Path,
    *,
    settings: dict[str, Any],
    metadata: dict[str, Any],
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
) -> AsrResult:
    configured_model = str(settings["model_path"] or "").strip()
    if not configured_model:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ASR model path is not configured",
        )
    model_path = Path(configured_model).expanduser()
    if not model_path.exists():
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ASR model path not found",
        )
    model_key = _asr_segment_cache_key(_asr_model_key(model_path), metadata)
    cached = _read_asr_cache(path, model_key)
    if cached is not None:
        text, segments = cached
        return AsrResult(
            status="completed",
            text=text,
            metadata={**metadata, "status": "cache_hit", "cache_hits": 1, "cache_misses": 0, "segments": segments},
        )
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg is None:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ffmpeg command not found",
        )
    if importlib.util.find_spec("faster_whisper") is None:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="faster_whisper not installed",
        )
    return _asr_with_faster_whisper(
        path,
        ffmpeg=ffmpeg,
        model_path=model_path,
        model_key=model_key,
        metadata=metadata,
        device=settings["device"],
        compute_type=settings["compute_type"],
        segment_start_seconds=segment_start_seconds,
        segment_duration_seconds=segment_duration_seconds,
    )


def _asr_media_openai_compatible(
    path: Path,
    *,
    settings: dict[str, Any],
    metadata: dict[str, Any],
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
) -> AsrResult:
    model = str(settings.get("model") or "").strip()
    base_url = str(settings.get("base_url") or "").strip().rstrip("/")
    if not model:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ASR model is not configured",
        )
    if not base_url:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ASR service base URL is not configured",
        )
    model_key = _asr_segment_cache_key(_asr_provider_model_key("openai_compatible", model), metadata)
    cached = _read_asr_cache(path, model_key)
    if cached is not None:
        text, segments = cached
        return AsrResult(
            status="completed",
            text=text,
            metadata={
                **metadata,
                "base_url": base_url,
                "status": "cache_hit",
                "cache_hits": 1,
                "cache_misses": 0,
                "segments": segments,
            },
        )
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg is None:
        return AsrResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "base_url": base_url, "status": "blocked_missing_dependency"},
            message="ffmpeg command not found",
        )
    return _asr_with_openai_compatible(
        path,
        ffmpeg=ffmpeg,
        model=model,
        base_url=base_url,
        model_key=model_key,
        metadata={**metadata, "base_url": base_url},
        max_duration_seconds=int(settings["max_duration_seconds"]),
        segment_start_seconds=segment_start_seconds,
        segment_duration_seconds=segment_duration_seconds,
    )


def _asr_settings() -> dict[str, Any]:
    defaults = {
        "enabled": True,
        "provider": "local_faster_whisper",
        "model": "",
        "base_url": "",
        "model_path": "",
        "max_duration_seconds": 3600,
        "device": "auto",
        "compute_type": "default",
    }
    try:
        from .settings import SettingsService

        service = SettingsService()
        return {
            "enabled": bool(service.resolve("acceleration.asr.enabled").raw_value),
            "provider": str(service.resolve("acceleration.asr.provider").raw_value or "local_faster_whisper"),
            "model": str(service.resolve("acceleration.asr.model").raw_value or ""),
            "base_url": str(service.resolve("acceleration.asr.base_url").raw_value or ""),
            "model_path": str(service.resolve("acceleration.asr.model_path").raw_value or ""),
            "max_duration_seconds": int(service.resolve("acceleration.asr.max_duration_seconds").raw_value or 3600),
            "device": str(service.resolve("acceleration.asr.device").raw_value or "auto"),
            "compute_type": str(service.resolve("acceleration.asr.compute_type").raw_value or "default"),
        }
    except Exception:
        return defaults


def _asr_with_faster_whisper(
    path: Path,
    *,
    ffmpeg: str,
    model_path: Path,
    model_key: str,
    metadata: dict[str, Any],
    device: str,
    compute_type: str,
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
) -> AsrResult:
    with tempfile.TemporaryDirectory(prefix="flux-kb-asr-") as temp_dir:
        audio_path = Path(temp_dir) / "audio.wav"
        extract_error = _extract_asr_audio(
            path,
            ffmpeg=ffmpeg,
            audio_path=audio_path,
            metadata=metadata,
            segment_start_seconds=segment_start_seconds,
            segment_duration_seconds=segment_duration_seconds,
        )
        if extract_error is not None:
            return extract_error
        try:
            faster_whisper = importlib.import_module("faster_whisper")
            model_kwargs: dict[str, Any] = {"local_files_only": True}
            if device != "auto":
                model_kwargs["device"] = device
            if compute_type != "default":
                model_kwargs["compute_type"] = compute_type
            model = faster_whisper.WhisperModel(str(model_path), **model_kwargs)
            segments_iter, info = model.transcribe(str(audio_path))
            parts: list[str] = []
            segment_count = 0
            for segment in segments_iter:
                text = str(getattr(segment, "text", "") or "").strip()
                if text:
                    parts.append(text)
                segment_count += 1
            redacted, _ = redact_text("\n".join(parts).strip())
            text = redacted.strip()
            _write_asr_cache(path, model_key, "faster_whisper", text, segment_count)
            result_metadata = {
                **metadata,
                "status": "completed",
                "cache_hits": 0,
                "cache_misses": 1,
                "segments": segment_count,
            }
            language = getattr(info, "language", None)
            if language:
                result_metadata["language"] = language
            return AsrResult(status="completed", text=text, metadata=result_metadata)
        except Exception as exc:
            return AsrResult(status="failed", metadata={**metadata, "status": "failed"}, message=str(exc))


def _asr_with_openai_compatible(
    path: Path,
    *,
    ffmpeg: str,
    model: str,
    base_url: str,
    model_key: str,
    metadata: dict[str, Any],
    max_duration_seconds: int,
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
) -> AsrResult:
    with tempfile.TemporaryDirectory(prefix="flux-kb-asr-") as temp_dir:
        audio_path = Path(temp_dir) / "audio.wav"
        extract_error = _extract_asr_audio(
            path,
            ffmpeg=ffmpeg,
            audio_path=audio_path,
            metadata=metadata,
            segment_start_seconds=segment_start_seconds,
            segment_duration_seconds=segment_duration_seconds,
        )
        if extract_error is not None:
            return extract_error
        try:
            timeout_seconds = int(segment_duration_seconds or max_duration_seconds)
            payload = _post_openai_compatible_asr(
                audio_path,
                base_url=base_url,
                model=model,
                timeout=max(ASR_FFMPEG_TIMEOUT_SECONDS, timeout_seconds + 60),
            )
        except _AsrServiceUnavailable as exc:
            return AsrResult(
                status="blocked_missing_dependency",
                metadata={**metadata, "status": "blocked_missing_dependency"},
                message=f"ASR service unavailable: {exc}",
            )
        except _AsrServiceError as exc:
            return AsrResult(status="failed", metadata={**metadata, "status": "failed"}, message=str(exc))
        text_value = payload.get("text")
        if not isinstance(text_value, str):
            return AsrResult(status="failed", metadata={**metadata, "status": "failed"}, message="ASR service response missing text")
        segments_value = payload.get("segments")
        segment_count = len(segments_value) if isinstance(segments_value, list) else 0
        redacted, _ = redact_text(text_value.strip())
        text = redacted.strip()
        _write_asr_cache(path, model_key, "openai_compatible_asr", text, segment_count)
        return AsrResult(
            status="completed",
            text=text,
            metadata={
                **metadata,
                "engine": "openai_compatible_asr",
                "status": "completed",
                "cache_hits": 0,
                "cache_misses": 1,
                "segments": segment_count,
            },
        )


def _extract_asr_audio(
    path: Path,
    *,
    ffmpeg: str,
    audio_path: Path,
    metadata: dict[str, Any],
    segment_start_seconds: float | None = None,
    segment_duration_seconds: float | None = None,
) -> AsrResult | None:
    command = [
        ffmpeg,
        "-y",
    ]
    if segment_start_seconds is not None and segment_start_seconds > 0:
        command.extend(["-ss", f"{segment_start_seconds:.3f}"])
    command.extend(
        [
        "-i",
        str(path),
        ]
    )
    if segment_duration_seconds is not None and segment_duration_seconds > 0:
        command.extend(["-t", f"{segment_duration_seconds:.3f}"])
    command.extend(
        [
        "-vn",
        "-ac",
        "1",
        "-ar",
        str(ASR_AUDIO_SAMPLE_RATE),
        "-f",
        "wav",
        str(audio_path),
        ]
    )
    timeout_seconds = int(segment_duration_seconds or metadata.get("duration_seconds") or ASR_FFMPEG_TIMEOUT_SECONDS)
    try:
        extract = run_no_window(
            command,
            text=True,
            capture_output=True,
            timeout=max(ASR_FFMPEG_TIMEOUT_SECONDS, timeout_seconds + 60),
            check=False,
        )
    except Exception as exc:  # pragma: no cover - environment-specific
        return AsrResult(status="failed", metadata={**metadata, "status": "failed"}, message=str(exc))
    if extract.returncode != 0:
        return AsrResult(
            status="failed",
            metadata={**metadata, "status": "failed"},
            message=extract.stderr.strip() or "ffmpeg failed",
        )
    return None


class _AsrServiceUnavailable(Exception):
    pass


class _AsrServiceError(Exception):
    pass


def _post_openai_compatible_asr(audio_path: Path, *, base_url: str, model: str, timeout: int) -> dict[str, Any]:
    boundary = hashlib.sha256(os.urandom(16)).hexdigest()
    body = _asr_multipart_body(
        boundary=boundary,
        fields={"model": model, "response_format": "json"},
        file_name="audio.wav",
        file_content=audio_path.read_bytes(),
    )
    request = Request(
        f"{base_url.rstrip('/')}/v1/audio/transcriptions",
        data=body,
        headers={"Content-Type": f"multipart/form-data; boundary={boundary}"},
        method="POST",
    )
    try:
        response = urlopen(request, timeout=timeout)
        raw = response.read(16 * 1024 * 1024)
    except HTTPError as exc:
        detail = _read_http_error_detail(exc)
        if exc.code == 503:
            raise _AsrServiceUnavailable(detail or "HTTP 503") from exc
        raise _AsrServiceError(detail or f"ASR service returned HTTP {exc.code}") from exc
    except (OSError, TimeoutError, URLError) as exc:
        reason = getattr(exc, "reason", None) or str(exc)
        raise _AsrServiceUnavailable(str(reason)) from exc
    try:
        payload = json.loads(raw.decode("utf-8"))
    except Exception as exc:
        raise _AsrServiceError("ASR service returned invalid JSON") from exc
    if not isinstance(payload, dict):
        raise _AsrServiceError("ASR service returned a non-object JSON response")
    return payload


def _asr_multipart_body(
    *,
    boundary: str,
    fields: dict[str, str],
    file_name: str,
    file_content: bytes,
) -> bytes:
    chunks: list[bytes] = []
    for name, value in fields.items():
        chunks.extend(
            [
                f"--{boundary}\r\n".encode("utf-8"),
                f'Content-Disposition: form-data; name="{name}"\r\n\r\n'.encode("utf-8"),
                str(value).encode("utf-8"),
                b"\r\n",
            ]
        )
    chunks.extend(
        [
            f"--{boundary}\r\n".encode("utf-8"),
            f'Content-Disposition: form-data; name="file"; filename="{file_name}"\r\n'.encode("utf-8"),
            b"Content-Type: audio/wav\r\n\r\n",
            file_content,
            b"\r\n",
            f"--{boundary}--\r\n".encode("utf-8"),
        ]
    )
    return b"".join(chunks)


def _read_http_error_detail(exc: HTTPError) -> str:
    try:
        raw = exc.read(1024 * 1024)
    except Exception:
        return ""
    text = raw.decode("utf-8", errors="replace").strip()
    if not text:
        return ""
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        return text
    if isinstance(payload, dict):
        detail = payload.get("detail")
        if isinstance(detail, str):
            return detail
    return text


def _asr_metadata(
    *,
    status: str,
    duration_seconds: float | None,
    max_duration_seconds: int,
    device: str,
    compute_type: str,
    provider: str,
    model: str,
    base_url: str,
) -> dict[str, Any]:
    engine = "openai_compatible_asr" if provider == "openai_compatible" else "faster_whisper"
    metadata: dict[str, Any] = {
        "engine": engine,
        "provider": provider,
        "status": status,
        "max_duration_seconds": max_duration_seconds,
        "device": device,
        "compute_type": compute_type,
        "cache_hits": 0,
        "cache_misses": 0,
        "segments": 0,
    }
    if model:
        metadata["model"] = model
    if base_url:
        metadata["base_url"] = base_url.rstrip("/")
    if duration_seconds is not None:
        metadata["duration_seconds"] = duration_seconds
    return metadata


def _media_duration_seconds(probed: dict[str, Any]) -> float | None:
    candidates: list[Any] = []
    if isinstance(probed.get("format"), dict):
        candidates.append(probed["format"].get("duration"))
    streams = probed.get("streams")
    if isinstance(streams, list):
        candidates.extend(stream.get("duration") for stream in streams if isinstance(stream, dict))
    for value in candidates:
        try:
            parsed = float(str(value).strip())
        except (TypeError, ValueError):
            continue
        if parsed >= 0:
            return parsed
    return None


def _media_has_audio_stream(probed: dict[str, Any]) -> bool | None:
    streams = probed.get("streams")
    if not isinstance(streams, list):
        return None
    return any(isinstance(stream, dict) and stream.get("codec_type") == "audio" for stream in streams)


def _sample_video_frames(path: Path, probed: dict[str, Any]) -> FrameSamplingResult:
    settings = _frame_sampling_settings()
    duration = _media_duration_seconds(probed)
    metadata = _frame_sampling_metadata(
        status="pending",
        duration_seconds=duration,
        max_duration_seconds=settings["max_duration_seconds"],
        frame_sample_count=settings["frame_sample_count"],
        scene_threshold=settings["scene_threshold"],
    )
    if not settings["enabled"]:
        return FrameSamplingResult(status="disabled", metadata={**metadata, "status": "disabled"})
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg is None:
        return FrameSamplingResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ffmpeg command not found",
        )
    transitions = _detect_scene_transitions(path, ffmpeg=ffmpeg, threshold=settings["scene_threshold"])
    if transitions is None:
        return FrameSamplingResult(status="failed", metadata={**metadata, "status": "failed"}, message="ffmpeg scene detection failed")
    selected = _selected_transition_frames(transitions, settings["frame_sample_count"])
    status = "completed"
    if not selected:
        if duration is None:
            return FrameSamplingResult(status="completed", metadata={**metadata, "status": "skipped_no_duration"})
        selected = [(round(duration / 2, 3), None)]
        status = "fallback_no_transition"
    chunks: list[AssetChunk] = []
    timestamps: list[float] = []
    scene_scores: list[float] = []
    thumbnail_hits = 0
    thumbnail_misses = 0
    vision_summary = {
        "engine": "local_vision",
        "provider": _vision_settings()["provider"],
        "model": _vision_settings()["model"],
        "prompt_schema": VISION_PROMPT_SCHEMA,
        "status": "disabled",
        "cache_hits": 0,
        "cache_misses": 0,
        "descriptions": 0,
        "blocked_dependency_count": 0,
    }
    for timestamp, score in selected:
        thumbnail = _thumbnail_for_frame(path, ffmpeg=ffmpeg, timestamp=timestamp)
        if thumbnail["status"] == "failed":
            return FrameSamplingResult(
                status="failed",
                metadata={**metadata, "status": "failed", "frame_count": len(timestamps)},
                message=str(thumbnail.get("message") or "thumbnail extraction failed"),
            )
        if thumbnail["cache_hit"]:
            thumbnail_hits += 1
        else:
            thumbnail_misses += 1
        timestamps.append(round(timestamp, 3))
        if score is not None:
            scene_scores.append(round(score, 6))
        vision = _vision_image(Path(thumbnail["path"]), source_label=f"{path.name}@{timestamp:.3f}s")
        if vision.status != "disabled":
            vision_summary["status"] = "completed" if vision.status == "completed" else str(vision.metadata.get("status") or vision.status)
            vision_summary["cache_hits"] += int(vision.metadata.get("cache_hits") or 0)
            vision_summary["cache_misses"] += int(vision.metadata.get("cache_misses") or 0)
            vision_summary["descriptions"] += int(vision.metadata.get("descriptions") or 0)
            vision_summary["blocked_dependency_count"] += int(vision.metadata.get("blocked_dependency_count") or 0)
        if vision.text:
            chunks.extend(_chunks_from_text(vision.text, f"{path.name} frame {timestamp:.3f}s", modality="vision"))
    final_metadata = {
        **metadata,
        "status": status,
        "timestamps": timestamps,
        "scene_scores": scene_scores,
        "frame_count": len(timestamps),
        "thumbnail_cache_hits": thumbnail_hits,
        "thumbnail_cache_misses": thumbnail_misses,
    }
    if vision_summary["status"] != "disabled":
        final_metadata["vision"] = vision_summary
    return FrameSamplingResult(status="completed", chunks=_reindexed_chunks(chunks), metadata=final_metadata)


def _sample_video_frames_at_timestamps(path: Path, probed: dict[str, Any], *, timestamps: list[float]) -> FrameSamplingResult:
    settings = _frame_sampling_settings()
    duration = _media_duration_seconds(probed)
    metadata = _frame_sampling_metadata(
        status="pending",
        duration_seconds=duration,
        max_duration_seconds=settings["max_duration_seconds"],
        frame_sample_count=settings["frame_sample_count"],
        scene_threshold=settings["scene_threshold"],
    )
    metadata["sampling_strategy"] = "staged_timestamps"
    if not settings["enabled"]:
        return FrameSamplingResult(status="disabled", metadata={**metadata, "status": "disabled"})
    if not timestamps:
        return FrameSamplingResult(status="completed", metadata={**metadata, "status": "skipped_no_timestamps"})
    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg is None:
        return FrameSamplingResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency"},
            message="ffmpeg command not found",
        )
    chunks: list[AssetChunk] = []
    thumbnail_hits = 0
    thumbnail_misses = 0
    sampled: list[float] = []
    vision_summary = {
        "engine": "local_vision",
        "provider": _vision_settings()["provider"],
        "model": _vision_settings()["model"],
        "prompt_schema": VISION_PROMPT_SCHEMA,
        "status": "disabled",
        "cache_hits": 0,
        "cache_misses": 0,
        "descriptions": 0,
        "blocked_dependency_count": 0,
    }
    for timestamp in timestamps:
        thumbnail = _thumbnail_for_frame(path, ffmpeg=ffmpeg, timestamp=timestamp)
        if thumbnail["status"] == "failed":
            return FrameSamplingResult(
                status="failed",
                metadata={**metadata, "status": "failed", "frame_count": len(sampled)},
                message=str(thumbnail.get("message") or "thumbnail extraction failed"),
            )
        if thumbnail["cache_hit"]:
            thumbnail_hits += 1
        else:
            thumbnail_misses += 1
        sampled.append(round(float(timestamp), 3))
        vision = _vision_image(Path(thumbnail["path"]), source_label=f"{path.name}@{timestamp:.3f}s")
        if vision.status != "disabled":
            vision_summary["status"] = "completed" if vision.status == "completed" else str(vision.metadata.get("status") or vision.status)
            vision_summary["cache_hits"] += int(vision.metadata.get("cache_hits") or 0)
            vision_summary["cache_misses"] += int(vision.metadata.get("cache_misses") or 0)
            vision_summary["descriptions"] += int(vision.metadata.get("descriptions") or 0)
            vision_summary["blocked_dependency_count"] += int(vision.metadata.get("blocked_dependency_count") or 0)
        if vision.text:
            chunks.extend(_chunks_from_text(vision.text, f"{path.name} frame {timestamp:.3f}s", modality="vision"))
    final_metadata = {
        **metadata,
        "status": "completed",
        "timestamps": sampled,
        "scene_scores": [],
        "frame_count": len(sampled),
        "thumbnail_cache_hits": thumbnail_hits,
        "thumbnail_cache_misses": thumbnail_misses,
    }
    if vision_summary["status"] != "disabled":
        final_metadata["vision"] = vision_summary
    return FrameSamplingResult(status="completed", chunks=_reindexed_chunks(chunks), metadata=final_metadata)


def _frame_sampling_settings() -> dict[str, Any]:
    defaults = {
        "enabled": True,
        "frame_sample_count": 3,
        "scene_threshold": 0.35,
        "max_duration_seconds": 1800,
    }
    try:
        from .settings import SettingsService

        service = SettingsService()
        return {
            "enabled": bool(service.resolve("acceleration.video.frame_sampling.enabled").raw_value),
            "frame_sample_count": int(service.resolve("acceleration.video.frame_sample_count").raw_value or 3),
            "scene_threshold": float(service.resolve("acceleration.video.scene_threshold").raw_value or 0.35),
            "max_duration_seconds": int(service.resolve("acceleration.video.frame_max_duration_seconds").raw_value or 1800),
        }
    except Exception:
        return defaults


def _frame_sampling_metadata(
    *,
    status: str,
    duration_seconds: float | None,
    max_duration_seconds: int,
    frame_sample_count: int,
    scene_threshold: float,
) -> dict[str, Any]:
    metadata: dict[str, Any] = {
        "engine": "ffmpeg",
        "status": status,
        "max_duration_seconds": max_duration_seconds,
        "frame_sample_count": frame_sample_count,
        "scene_threshold": scene_threshold,
        "frame_count": 0,
        "thumbnail_cache_hits": 0,
        "thumbnail_cache_misses": 0,
    }
    if duration_seconds is not None:
        metadata["duration_seconds"] = duration_seconds
    return metadata


def _detect_scene_transitions(path: Path, *, ffmpeg: str, threshold: float) -> list[tuple[float, float]] | None:
    threshold_text = f"{threshold:g}"
    try:
        result = run_no_window(
            [
                ffmpeg,
                "-hide_banner",
                "-i",
                str(path),
                "-vf",
                f"select='gt(scene,{threshold_text})',metadata=print",
                "-an",
                "-f",
                "null",
                "-",
            ],
            text=True,
            capture_output=True,
            timeout=FRAME_SAMPLING_TIMEOUT_SECONDS,
            check=False,
        )
    except Exception:
        return None
    if result.returncode != 0:
        return None
    return _parse_scene_transitions("\n".join([str(result.stdout or ""), str(result.stderr or "")]))


def _parse_scene_transitions(text: str) -> list[tuple[float, float]]:
    transitions: list[tuple[float, float]] = []
    pending_timestamp: float | None = None
    for line in text.splitlines():
        timestamp_match = re.search(r"pts_time[:=]\s*([0-9]+(?:\.[0-9]+)?)", line)
        score_match = re.search(r"(?:lavfi\.)?scene_score[=:]\s*([0-9]+(?:\.[0-9]+)?)", line)
        try:
            timestamp = float(timestamp_match.group(1)) if timestamp_match else pending_timestamp
            if timestamp_match:
                pending_timestamp = timestamp
            score = float(score_match.group(1)) if score_match else None
        except ValueError:
            continue
        if timestamp is None or score is None:
            continue
        transitions.append((timestamp, score))
        pending_timestamp = None
    return transitions


def _selected_transition_frames(transitions: list[tuple[float, float]], limit: int) -> list[tuple[float, float]]:
    strongest = sorted(transitions, key=lambda item: (-item[1], item[0]))[: max(1, limit)]
    return sorted(strongest, key=lambda item: item[0])


def _thumbnail_for_frame(path: Path, *, ffmpeg: str, timestamp: float) -> dict[str, Any]:
    cache_file = _thumbnail_cache_file(path, timestamp=timestamp)
    if cache_file.exists():
        return {"status": "completed", "path": str(cache_file), "cache_hit": True}
    try:
        cache_file.parent.mkdir(parents=True, exist_ok=True)
        result = run_no_window(
            [
                ffmpeg,
                "-y",
                "-ss",
                f"{timestamp:.3f}",
                "-i",
                str(path),
                "-frames:v",
                "1",
                "-f",
                "image2",
                str(cache_file),
            ],
            text=True,
            capture_output=True,
            timeout=FRAME_SAMPLING_TIMEOUT_SECONDS,
            check=False,
        )
    except Exception as exc:
        return {"status": "failed", "path": str(cache_file), "cache_hit": False, "message": str(exc)}
    if result.returncode != 0:
        return {
            "status": "failed",
            "path": str(cache_file),
            "cache_hit": False,
            "message": result.stderr.strip() or "ffmpeg thumbnail extraction failed",
        }
    if not cache_file.exists():
        return {"status": "failed", "path": str(cache_file), "cache_hit": False, "message": "thumbnail file was not created"}
    return {"status": "completed", "path": str(cache_file), "cache_hit": False}


def _thumbnail_cache_file(path: Path, *, timestamp: float) -> Path:
    source_hash = _sha256_file(path)
    timestamp_key = f"{timestamp:.3f}"
    key = hashlib.sha256(f"{THUMBNAIL_CACHE_SCHEMA}:{source_hash}:{timestamp_key}".encode("utf-8")).hexdigest()
    return Path(resolve_cache_layout()["directories"]["thumbnails"]) / f"{key}.png"


def _vision_summary_from_frame_sampling(frame_sampling: FrameSamplingResult) -> dict[str, Any] | None:
    vision = frame_sampling.metadata.get("vision") if isinstance(frame_sampling.metadata, dict) else None
    return dict(vision) if isinstance(vision, dict) else None


def _read_asr_cache(path: Path, model_key: str) -> tuple[str, int] | None:
    try:
        cache_file = _asr_cache_file(path, model_key)
        payload = json.loads(cache_file.read_text(encoding="utf-8"))
    except Exception:
        return None
    if payload.get("schema") != ASR_CACHE_SCHEMA:
        return None
    if payload.get("source_hash") != _sha256_file(path):
        return None
    if payload.get("model_key") != model_key:
        return None
    text = payload.get("text")
    segments = payload.get("segments")
    if not isinstance(text, str):
        return None
    try:
        segment_count = int(segments or 0)
    except (TypeError, ValueError):
        segment_count = 0
    return text, segment_count


def _write_asr_cache(path: Path, model_key: str, engine: str, text: str, segments: int) -> None:
    try:
        cache_file = _asr_cache_file(path, model_key)
        cache_file.parent.mkdir(parents=True, exist_ok=True)
        payload = {
            "schema": ASR_CACHE_SCHEMA,
            "source_hash": _sha256_file(path),
            "model_key": model_key,
            "engine": engine,
            "segments": int(segments),
            "text": text,
        }
        cache_file.write_text(json.dumps(payload, sort_keys=True), encoding="utf-8")
    except Exception:
        return


def _asr_cache_file(path: Path, model_key: str) -> Path:
    source_hash = _sha256_file(path)
    key = hashlib.sha256(f"{ASR_CACHE_SCHEMA}:{model_key}:{source_hash}".encode("utf-8")).hexdigest()
    return Path(resolve_cache_layout()["directories"]["asr"]) / f"{key}.json"


def _asr_model_key(model_path: Path) -> str:
    try:
        resolved = str(model_path.resolve())
    except Exception:
        resolved = str(model_path)
    return _asr_provider_model_key("local_faster_whisper", resolved)


def _asr_provider_model_key(provider: str, model: str) -> str:
    return f"{provider}:{model}"


def _asr_segment_cache_key(model_key: str, metadata: dict[str, Any]) -> str:
    if "segment_index" not in metadata:
        return model_key
    start = metadata.get("segment_start_seconds")
    duration = metadata.get("segment_duration_seconds")
    return f"{model_key}:segment:{metadata.get('segment_index')}:{start}:{duration}"


def _chunks_from_text(text: str, title: str, *, modality: str = "text") -> tuple[AssetChunk, ...]:
    cleaned = text.strip()
    if not cleaned:
        return ()
    redacted, _ = redact_text(cleaned)
    chunks: list[AssetChunk] = []
    for index, start in enumerate(range(0, len(redacted), 4000)):
        body = redacted[start : start + 4000].strip()
        if body:
            chunks.append(
                AssetChunk(
                    chunk_index=index,
                    title=title,
                    body=body,
                    modality=modality,
                    locator=f"char:{start}-{start + len(body)}",
                    token_estimate=max(1, len(body.split())),
                )
            )
    return tuple(chunks)


def _reindexed_chunks(*groups: Iterable[AssetChunk]) -> tuple[AssetChunk, ...]:
    chunks: list[AssetChunk] = []
    for group in groups:
        for chunk in group:
            chunks.append(replace(chunk, chunk_index=len(chunks)))
    return tuple(chunks)


def _offset_chunks(
    chunks: Iterable[AssetChunk],
    offset: int,
    *,
    locator_prefix: str | None = None,
    metadata: dict[str, Any] | None = None,
) -> tuple[AssetChunk, ...]:
    adjusted: list[AssetChunk] = []
    for chunk in chunks:
        locator = chunk.locator
        if locator_prefix:
            locator = f"{locator_prefix}:{locator}" if locator else locator_prefix
        chunk_metadata = dict(chunk.metadata or {})
        if metadata:
            chunk_metadata.update(metadata)
        adjusted.append(
            replace(
                chunk,
                chunk_index=offset + int(chunk.chunk_index),
                locator=locator,
                metadata=chunk_metadata,
            )
        )
    return tuple(adjusted)


def _decorative_image_metadata(metadata: dict[str, Any]) -> dict[str, Any] | None:
    width = _int_or_none(metadata.get("width"))
    height = _int_or_none(metadata.get("height"))
    if width is not None and height is not None and (width <= 1 or height <= 1):
        return {"status": "skipped", "reason": "tiny_spacer"}
    size_bytes = _int_or_none(metadata.get("size_bytes"))
    if (
        width is not None
        and height is not None
        and size_bytes is not None
        and metadata.get("has_transparency") is True
        and width <= DECORATIVE_ICON_MAX_DIMENSION
        and height <= DECORATIVE_ICON_MAX_DIMENSION
        and size_bytes <= DECORATIVE_ICON_MAX_BYTES
    ):
        return {"status": "skipped", "reason": "small_icon"}
    return None


def _vision_image(path: Path, *, source_label: str) -> VisionResult:
    settings = _vision_settings()
    metadata = _vision_metadata(
        status="pending",
        provider=settings["provider"],
        model=settings["model"],
    )
    if not settings["enabled"]:
        return VisionResult(status="disabled", metadata={**metadata, "status": "disabled"})
    dimensions = _image_dimensions(path)
    if dimensions:
        pixels = dimensions[0] * dimensions[1]
        metadata["width"] = dimensions[0]
        metadata["height"] = dimensions[1]
        metadata["pixels"] = pixels
        if pixels > settings["max_image_pixels"]:
            return VisionResult(status="completed", metadata={**metadata, "status": "skipped_pixel_cap"})
    if not settings["local_inference_enabled"]:
        return VisionResult(
            status="blocked_config",
            metadata={**metadata, "status": "blocked_config", "blocked_dependency_count": 1},
            message="local inference is disabled",
        )
    if settings["provider"] != "ollama":
        return VisionResult(
            status="blocked_config",
            metadata={**metadata, "status": "blocked_config", "blocked_dependency_count": 1},
            message=f"vision provider {settings['provider']} is not implemented for vision enrichment in this build",
        )
    if not settings["model"]:
        return VisionResult(
            status="blocked_missing_dependency",
            metadata={**metadata, "status": "blocked_missing_dependency", "blocked_dependency_count": 1},
            message="vision model is not configured",
        )
    try:
        base_url = validate_local_model_base_url(settings["base_url"])
    except ValueError as exc:
        return VisionResult(
            status="blocked_config",
            metadata={**metadata, "status": "blocked_config", "blocked_dependency_count": 1},
            message=str(exc),
        )
    cached = _read_vision_cache(path, provider=settings["provider"], model=settings["model"])
    if cached is not None:
        return VisionResult(
            status="completed",
            text=cached,
            metadata={**metadata, "status": "cache_hit", "cache_hits": 1, "cache_misses": 0, "descriptions": 1 if cached else 0},
        )
    return _vision_with_ollama_compatible(
        path,
        source_label=source_label,
        provider=settings["provider"],
        base_url=base_url,
        model=settings["model"],
        keep_alive=settings["keep_alive"],
        timeout_seconds=settings["timeout_seconds"],
        metadata=metadata,
    )


def _vision_settings() -> dict[str, Any]:
    defaults = {
        "enabled": True,
        "model": "qwen3-vl:8b",
        "max_image_pixels": 4_096_000,
        "local_inference_enabled": True,
        "provider": "ollama",
        "base_url": "http://127.0.0.1:11434",
        "keep_alive": "2m",
        "timeout_seconds": 1,
    }
    try:
        from .settings import SettingsService

        service = SettingsService()
        return {
            "enabled": bool(service.resolve("acceleration.vision.enabled").raw_value),
            "model": str(service.resolve("acceleration.vision.model").raw_value or ""),
            "max_image_pixels": int(service.resolve("acceleration.vision.max_image_pixels").raw_value or 4_096_000),
            "local_inference_enabled": bool(service.resolve("acceleration.local_inference.enabled").raw_value),
            "provider": str(service.resolve("acceleration.local_inference.provider").raw_value or "ollama"),
            "base_url": str(service.resolve("acceleration.local_inference.base_url").raw_value or "http://127.0.0.1:11434"),
            "keep_alive": str(service.resolve("acceleration.local_inference.keep_alive").raw_value or ""),
            "timeout_seconds": int(service.resolve("acceleration.local_inference.probe_timeout_seconds").raw_value or 1),
        }
    except Exception:
        return defaults


def _vision_metadata(*, status: str, provider: str, model: str) -> dict[str, Any]:
    return {
        "engine": "local_vision",
        "provider": provider,
        "model": model,
        "prompt_schema": VISION_PROMPT_SCHEMA,
        "status": status,
        "cache_hits": 0,
        "cache_misses": 0,
        "descriptions": 0,
        "blocked_dependency_count": 0,
    }


def _vision_with_ollama_compatible(
    path: Path,
    *,
    source_label: str,
    provider: str,
    base_url: str,
    model: str,
    keep_alive: str,
    timeout_seconds: int,
    metadata: dict[str, Any],
) -> VisionResult:
    request_attempted = False
    try:
        image_bytes, submission_metadata = _vision_request_image_bytes(path)
        image_b64 = base64.b64encode(image_bytes).decode("ascii")
        payload = {
            "model": model,
            "stream": False,
            "prompt": (
                "Describe this image for a private local knowledge index. "
                "Be concise and factual. If it is a diagram, name the diagram type, visible title, "
                "key labels, entities, and relationships. Answer directly in plain text. "
                "Do not infer private identities."
            ),
            "images": [image_b64],
            "options": {"num_predict": VISION_NUM_PREDICT},
        }
        if keep_alive:
            payload["keep_alive"] = keep_alive
        request = Request(
            f"{base_url}/api/generate",
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        request_attempted = True
        response = urlopen(request, timeout=max(1, timeout_seconds, VISION_TIMEOUT_SECONDS))
        raw = response.read(2 * 1024 * 1024).decode("utf-8", errors="replace")
        decoded = json.loads(raw or "{}")
        text, response_field = _vision_response_text(decoded)
        redacted, _ = redact_text(text)
        clean = _clean_vision_caption_text(redacted)
        if clean:
            _write_vision_cache(path, provider=provider, model=model, text=clean)
        response_metadata = _vision_response_metadata(decoded)
        if response_field and response_field != "response":
            response_metadata["fallback_field"] = response_field
        return VisionResult(
            status="completed",
            text=clean,
            metadata={
                **metadata,
                **submission_metadata,
                **response_metadata,
                "status": "completed",
                "cache_hits": 0,
                "cache_misses": 1,
                "descriptions": 1 if clean else 0,
                "source_label": source_label,
            },
        )
    except Exception as exc:
        message = str(exc)
        failed_metadata = {**metadata, "status": "failed", "error": message[:500]}
        if request_attempted:
            failed_metadata["cache_misses"] = 1
        return VisionResult(status="failed", metadata=failed_metadata, message=message)


def _vision_request_image_bytes(path: Path) -> tuple[bytes, dict[str, Any]]:
    try:
        from PIL import Image

        with Image.open(path) as image:
            width, height = image.size
            metadata: dict[str, Any] = {
                "input_width": width,
                "input_height": height,
                "submitted_max_edge": VISION_REQUEST_MAX_EDGE,
            }
            if max(width, height) <= VISION_REQUEST_MAX_EDGE:
                data = path.read_bytes()
                metadata.update(
                    {
                        "submitted_width": width,
                        "submitted_height": height,
                        "submitted_bytes": len(data),
                        "submitted_resize": "original",
                    }
                )
                return data, metadata
            resized = image.convert("RGB")
            resized.thumbnail((VISION_REQUEST_MAX_EDGE, VISION_REQUEST_MAX_EDGE), Image.Resampling.LANCZOS)
            buffer = BytesIO()
            resized.save(buffer, format="JPEG", quality=92, optimize=True)
            data = buffer.getvalue()
            metadata.update(
                {
                    "submitted_width": resized.width,
                    "submitted_height": resized.height,
                    "submitted_bytes": len(data),
                    "submitted_resize": "scaled",
                }
            )
            return data, metadata
    except Exception:
        data = path.read_bytes()
        return data, {"submitted_resize": "unavailable", "submitted_bytes": len(data)}


def _vision_response_text(payload: Any) -> tuple[str, str]:
    if isinstance(payload, dict):
        response = payload.get("response")
        if isinstance(response, str):
            if response.strip():
                return response, "response"
        message = payload.get("message")
        if isinstance(message, dict) and isinstance(message.get("content"), str):
            content = message["content"]
            if content.strip():
                return content, "message.content"
        choices = payload.get("choices")
        if isinstance(choices, list) and choices:
            first = choices[0]
            if isinstance(first, dict):
                if isinstance(first.get("text"), str):
                    text = first["text"]
                    if text.strip():
                        return text, "choices[0].text"
                nested = first.get("message")
                if isinstance(nested, dict) and isinstance(nested.get("content"), str):
                    content = nested["content"]
                    if content.strip():
                        return content, "choices[0].message.content"
        thinking = payload.get("thinking")
        if isinstance(thinking, str) and thinking.strip():
            return _strip_thinking_markup(thinking), "thinking"
    return "", ""


def _strip_thinking_markup(text: str) -> str:
    stripped = text.strip()
    stripped = re.sub(r"^\s*<think>\s*", "", stripped, flags=re.IGNORECASE)
    stripped = re.sub(r"\s*</think>\s*$", "", stripped, flags=re.IGNORECASE)
    return stripped.strip()


def _clean_vision_caption_text(text: str) -> str:
    cleaned = _strip_thinking_markup(text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    for pattern in (
        r"\b(?:The image is|This image is|This is|It is|It's|Diagram Type:|Entity[- ]Relationship Diagram)\b",
    ):
        match = re.search(pattern, cleaned, flags=re.IGNORECASE)
        if match:
            cleaned = cleaned[match.start() :]
            break
    cleaned = re.split(
        r"\b(?:Need to|Wait,|Check the user instructions|The user wants|They want)\b",
        cleaned,
        maxsplit=1,
        flags=re.IGNORECASE,
    )[0].strip()
    return cleaned


def _vision_response_metadata(payload: Any) -> dict[str, Any]:
    if not isinstance(payload, dict):
        return {}
    metadata: dict[str, Any] = {}
    for key in ("done_reason", "prompt_eval_count", "eval_count"):
        value = payload.get(key)
        if isinstance(value, (int, str)):
            metadata[key] = value
    return metadata


def _read_vision_cache(path: Path, *, provider: str, model: str) -> str | None:
    try:
        cache_file = _vision_cache_file(path, provider=provider, model=model)
        payload = json.loads(cache_file.read_text(encoding="utf-8"))
    except Exception:
        return None
    if payload.get("schema") != VISION_CACHE_SCHEMA:
        return None
    if payload.get("source_hash") != _sha256_file(path):
        return None
    if payload.get("provider") != provider or payload.get("model") != model:
        return None
    if payload.get("prompt_schema") != VISION_PROMPT_SCHEMA:
        return None
    text = payload.get("text")
    return text if isinstance(text, str) else None


def _write_vision_cache(path: Path, *, provider: str, model: str, text: str) -> None:
    try:
        cache_file = _vision_cache_file(path, provider=provider, model=model)
        cache_file.parent.mkdir(parents=True, exist_ok=True)
        payload = {
            "schema": VISION_CACHE_SCHEMA,
            "source_hash": _sha256_file(path),
            "provider": provider,
            "model": model,
            "prompt_schema": VISION_PROMPT_SCHEMA,
            "text": text,
        }
        cache_file.write_text(json.dumps(payload, sort_keys=True), encoding="utf-8")
    except Exception:
        return


def _vision_cache_file(path: Path, *, provider: str, model: str) -> Path:
    source_hash = _sha256_file(path)
    key = hashlib.sha256(f"{VISION_CACHE_SCHEMA}:{provider}:{model}:{VISION_PROMPT_SCHEMA}:{source_hash}".encode("utf-8")).hexdigest()
    return Path(resolve_cache_layout()["directories"]["vision"]) / f"{key}.json"


def _int_or_none(value: Any) -> int | None:
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _optional_int(value: Any) -> int | None:
    return _int_or_none(value)


def _float_or_none(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _image_dimensions(path: Path) -> tuple[int, int] | None:
    try:
        from PIL import Image

        with Image.open(path) as image:
            return image.size
    except Exception:
        return _png_dimensions(path)


def _image_has_transparency(path: Path) -> bool | None:
    try:
        from PIL import Image

        with Image.open(path) as image:
            if image.mode in {"RGBA", "LA"}:
                extrema = image.getextrema()
                if extrema and isinstance(extrema[-1], tuple):
                    return int(extrema[-1][0]) < 255
            return bool(image.info.get("transparency"))
    except Exception:
        return None


def _png_dimensions(path: Path) -> tuple[int, int] | None:
    try:
        with path.open("rb") as handle:
            header = handle.read(24)
        if header[:8] != b"\x89PNG\r\n\x1a\n":
            return None
        return struct.unpack(">II", header[16:24])
    except Exception:
        return None


def _ocr_image(path: Path) -> OcrResult:
    cached = _read_ocr_cache(path)
    if cached is not None:
        return OcrResult(
            status="completed",
            text=cached,
            metadata={
                "engine": "tesseract",
                "status": "cache_hit",
                "cache_hits": 1,
                "cache_misses": 0,
            },
        )
    tesseract = shutil.which("tesseract")
    if tesseract is None:
        return OcrResult(
            status="blocked_missing_dependency",
            metadata={
                "engine": "tesseract",
                "status": "blocked_missing_dependency",
                "cache_hits": 0,
                "cache_misses": 0,
            },
            message="tesseract command not found",
        )
    return _ocr_image_with_tesseract(path, tesseract)


def _ocr_pdf(path: Path, *, page_count: int) -> OcrResult:
    return _ocr_pdf_pages(path, page_start=1, page_end=page_count, page_count=page_count)


def _ocr_pdf_pages(
    path: Path,
    *,
    page_start: int,
    page_end: int,
    page_count: int,
    page_numbers: list[int] | None = None,
) -> OcrResult:
    page_start = max(1, page_start)
    page_end = min(page_count, max(page_start, page_end))
    pages = page_numbers or list(range(page_start, page_end + 1))
    base_metadata: dict[str, Any] = {
        "engine": "tesseract",
        "renderer": "pdftoppm",
        "page_count": page_count,
        "page_start": page_start,
        "page_end": page_end,
        "pages": pages,
        "pages_attempted": 0,
        "cache_hits": 0,
        "cache_misses": 0,
    }
    renderer = shutil.which("pdftoppm")
    if renderer is None:
        return OcrResult(
            status="blocked_missing_dependency",
            metadata={**base_metadata, "status": "blocked_missing_dependency"},
            message="pdftoppm command not found",
        )
    tesseract = shutil.which("tesseract")
    if tesseract is None:
        return OcrResult(
            status="blocked_missing_dependency",
            metadata={**base_metadata, "status": "blocked_missing_dependency"},
            message="tesseract command not found",
        )
    parts: list[str] = []
    with tempfile.TemporaryDirectory(prefix="flux-kb-pdf-ocr-") as temp_dir:
        temp_root = Path(temp_dir)
        for page_index, page_number in enumerate(pages, start=1):
            output_prefix = temp_root / "page"
            try:
                render = run_no_window(
                    [
                        renderer,
                        "-r",
                        str(OCR_PDF_DPI),
                        "-scale-to",
                        str(OCR_MAX_IMAGE_EDGE),
                        "-png",
                        "-f",
                        str(page_number),
                        "-l",
                        str(page_number),
                        str(path),
                        str(output_prefix),
                    ],
                    text=True,
                    capture_output=True,
                    timeout=OCR_TIMEOUT_SECONDS,
                    check=False,
                )
            except subprocess.TimeoutExpired as exc:
                return OcrResult(
                    status="blocked_missing_dependency",
                    metadata={**base_metadata, "pages_attempted": page_index - 1, "status": "blocked_timeout"},
                    message=str(exc),
                )
            except Exception as exc:  # pragma: no cover - environment-specific
                return OcrResult(
                    status="failed",
                    metadata={**base_metadata, "pages_attempted": page_index - 1, "status": "failed"},
                    message=str(exc),
                )
            if render.returncode != 0:
                return OcrResult(
                    status="failed",
                    metadata={**base_metadata, "pages_attempted": page_index - 1, "status": "failed"},
                    message=render.stderr.strip() or "pdftoppm failed",
                )
            rendered_path = output_prefix.with_name(f"{output_prefix.name}-{page_number}.png")
            page_ocr = _ocr_image_with_tesseract(rendered_path, tesseract)
            base_metadata["pages_attempted"] = page_index
            base_metadata["cache_hits"] += int(page_ocr.metadata.get("cache_hits") or 0)
            base_metadata["cache_misses"] += int(page_ocr.metadata.get("cache_misses") or 0)
            if page_ocr.status == "failed":
                return OcrResult(
                    status="failed",
                    metadata={**base_metadata, "status": "failed"},
                    message=page_ocr.message,
                )
            if page_ocr.status == "blocked_missing_dependency":
                return OcrResult(
                    status="blocked_missing_dependency",
                    metadata={**base_metadata, "status": page_ocr.metadata.get("status") or "blocked_missing_dependency"},
                    message=page_ocr.message,
                )
            if page_ocr.text:
                parts.append(page_ocr.text)
    return OcrResult(status="completed", text="\n".join(parts), metadata={**base_metadata, "status": "completed"})


def _ocr_image_with_tesseract(path: Path, tesseract: str, *, cache_path: Path | None = None) -> OcrResult:
    source_path = cache_path or path
    cached = _read_ocr_cache(source_path)
    if cached is not None:
        return OcrResult(
            status="completed",
            text=cached,
            metadata={
                "engine": "tesseract",
                "status": "cache_hit",
                "cache_hits": 1,
                "cache_misses": 0,
                "preprocess": {"status": "cache_hit", "max_edge": OCR_MAX_IMAGE_EDGE},
            },
        )
    preprocess = _ocr_preprocess_metadata(path)
    tesseract_input = path
    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None
    if preprocess.get("status") == "needs_scaling":
        temp_dir_obj = tempfile.TemporaryDirectory(prefix="flux-kb-ocr-image-")
        temp_root = Path(temp_dir_obj.name)
        tesseract_input, preprocess = _scaled_ocr_input(path, temp_root=temp_root, base_metadata=preprocess)
    try:
        result = run_no_window(
            [tesseract, str(tesseract_input), "stdout"],
            text=True,
            capture_output=True,
            timeout=OCR_TIMEOUT_SECONDS,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        return OcrResult(
            status="blocked_missing_dependency",
            metadata={
                "engine": "tesseract",
                "status": "blocked_timeout",
                "cache_hits": 0,
                "cache_misses": 1,
                "preprocess": preprocess,
            },
            message=str(exc),
        )
    except Exception as exc:
        return OcrResult(
            status="failed",
            metadata={
                "engine": "tesseract",
                "status": "failed",
                "cache_hits": 0,
                "cache_misses": 1,
                "preprocess": preprocess,
            },
            message=str(exc),
        )
    finally:
        if temp_dir_obj is not None:
            temp_dir_obj.cleanup()
    if result.returncode != 0:
        return OcrResult(
            status="completed",
            metadata={
                "engine": "tesseract",
                "status": "failed",
                "cache_hits": 0,
                "cache_misses": 1,
                "preprocess": preprocess,
            },
            message=result.stderr.strip() or "tesseract failed",
        )
    redacted, _ = redact_text(result.stdout.strip())
    text = redacted.strip()
    _write_ocr_cache(source_path, text)
    return OcrResult(
        status="completed",
        text=text,
        metadata={
            "engine": "tesseract",
            "status": "completed",
            "cache_hits": 0,
            "cache_misses": 1,
            "preprocess": preprocess,
        },
    )


def _ocr_preprocess_metadata(path: Path) -> dict[str, Any]:
    width = None
    height = None
    dimensions = _image_dimensions(path)
    if dimensions is not None:
        width, height = dimensions
    metadata: dict[str, Any] = {
        "status": "original",
        "max_edge": OCR_MAX_IMAGE_EDGE,
    }
    if width is not None and height is not None:
        metadata.update(
            {
                "input_width": width,
                "input_height": height,
                "output_width": width,
                "output_height": height,
            }
        )
        if max(width, height) > OCR_MAX_IMAGE_EDGE:
            metadata["status"] = "needs_scaling"
    return metadata


def _scaled_ocr_input(path: Path, *, temp_root: Path, base_metadata: dict[str, Any]) -> tuple[Path, dict[str, Any]]:
    try:
        from PIL import Image

        output_path = temp_root / "ocr-input.png"
        with Image.open(path) as image:
            image.thumbnail((OCR_MAX_IMAGE_EDGE, OCR_MAX_IMAGE_EDGE), Image.Resampling.LANCZOS)
            prepared = image.convert("RGB")
            prepared.save(output_path)
            output_width, output_height = prepared.size
        return output_path, {
            **base_metadata,
            "status": "scaled",
            "output_width": output_width,
            "output_height": output_height,
        }
    except Exception as exc:
        return path, {**base_metadata, "status": "scale_failed", "message": str(exc)[:200]}


def _read_ocr_cache(path: Path) -> str | None:
    try:
        cache_file = _ocr_cache_file(path)
        payload = json.loads(cache_file.read_text(encoding="utf-8"))
    except Exception:
        return None
    if payload.get("schema") != OCR_CACHE_SCHEMA:
        return None
    if payload.get("source_hash") != _sha256_file(path):
        return None
    text = payload.get("text")
    return text if isinstance(text, str) else None


def _write_ocr_cache(path: Path, text: str) -> None:
    try:
        cache_file = _ocr_cache_file(path)
        cache_file.parent.mkdir(parents=True, exist_ok=True)
        payload = {
            "schema": OCR_CACHE_SCHEMA,
            "source_hash": _sha256_file(path),
            "engine": "tesseract",
            "text": text,
        }
        cache_file.write_text(json.dumps(payload, sort_keys=True), encoding="utf-8")
    except Exception:
        return


def _ocr_cache_file(path: Path) -> Path:
    source_hash = _sha256_file(path)
    key = hashlib.sha256(f"{OCR_CACHE_SCHEMA}:tesseract:{source_hash}".encode("utf-8")).hexdigest()
    return Path(resolve_cache_layout()["directories"]["ocr"]) / f"{key}.json"


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _read_sidecar_transcript(path: Path) -> str | None:
    for suffix in (".txt", ".md", ".vtt", ".srt"):
        sidecar = path.with_suffix(path.suffix + suffix)
        if sidecar.exists():
            return read_text_with_bom(sidecar).strip()
    return None


def _module_check(module_name: str) -> dict[str, Any]:
    ok = importlib.util.find_spec(module_name) is not None
    return {"ok": ok, "message": "available" if ok else f"{module_name} not installed"}


def _tool_check(command: str) -> dict[str, Any]:
    path = shutil.which(command)
    return {"ok": path is not None, "message": path or f"{command} command not found"}


def _first_tool_check(label: str, commands: tuple[str, ...]) -> dict[str, Any]:
    for command in commands:
        path = shutil.which(command)
        if path is not None:
            return {"ok": True, "message": path}
    return {"ok": False, "message": f"{label} command not found"}


def _word_com_check() -> dict[str, Any]:
    if os.name != "nt":
        return {"ok": False, "message": "Windows Word COM is only available on Windows"}
    ok = importlib.util.find_spec("win32com") is not None
    return {"ok": ok, "message": "pywin32 available" if ok else "pywin32 not installed"}


def _excel_com_check() -> dict[str, Any]:
    if os.name != "nt":
        return {"ok": False, "message": "Windows Excel COM is only available on Windows"}
    ok = importlib.util.find_spec("win32com") is not None
    return {"ok": ok, "message": "pywin32 available" if ok else "pywin32 not installed"}


def _powerpoint_com_check() -> dict[str, Any]:
    if os.name != "nt":
        return {"ok": False, "message": "Windows PowerPoint COM is only available on Windows"}
    ok = importlib.util.find_spec("win32com") is not None
    return {"ok": ok, "message": "pywin32 available" if ok else "pywin32 not installed"}
